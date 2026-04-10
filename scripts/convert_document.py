#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
独立的文档转换脚本
- 将 Word、Excel、PowerPoint 和 PDF 文件转换为 Markdown 格式
- 将 Markdown 文件转换为 Word (.docx) 格式

依赖：
- python-docx: pip install python-docx
- openpyxl: pip install openpyxl
- python-pptx: pip install python-pptx
- pdfplumber: pip install pdfplumber
- Node.js: 用于 Markdown 转 DOCX（需要单独安装）
"""

import sys
import os
import json
import importlib
import logging
import re
import subprocess
import shutil
import io
import struct
import hashlib
import xml.etree.ElementTree as ET

SUPPORTED_EXTENSIONS = ['.docx', '.xlsx', '.pptx', '.pdf', '.md']
MAX_FILE_SIZE_BYTES = 100 * 1024 * 1024
NODE_CONVERT_TIMEOUT_SECONDS = 120
NODE_SHARED_HOME_ENV = "BRUCE_DOC_CONVERTER_NODE_HOME"
GENERATED_OUTPUT_DIR_NAMES = {"Markdown", "Word"}
DOCX_XML_NAMESPACES = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
DOCX_W_NS = DOCX_XML_NAMESPACES['w']

# 图片提取相关常量
IMAGE_OUTPUT_DIR_NAME = "images"
MIN_IMAGE_DIMENSION_PX = 20          # 小于此像素的图片视为装饰性
MAX_ASPECT_RATIO = 10.0              # 宽高比超过此值视为装饰线条
MIN_IMAGE_DATA_BYTES = 500           # 数据量低于此值视为纯色/透明占位
PPTX_BACKGROUND_COVERAGE_RATIO = 0.9 # 覆盖幻灯片面积超过此比例视为背景图

# OOXML 图片相关命名空间
OOXML_IMAGE_NAMESPACES = {
    'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
    'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
    'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006',
    'adec': 'http://schemas.microsoft.com/office/drawing/2017/decorative',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

# 图片格式文件头魔数
_IMAGE_SIGNATURES = {
    b'\x89PNG\r\n\x1a\n': 'png',
    b'\xff\xd8\xff': 'jpeg',
    b'GIF87a': 'gif',
    b'GIF89a': 'gif',
    b'BM': 'bmp',
    b'II\x2a\x00': 'tiff',
    b'MM\x00\x2a': 'tiff',
    b'\xd7\xcd\xc6\x9a': 'wmf',
    b'\x01\x00\x00\x00': 'emf',
}

logger = logging.getLogger(__name__)

_RE_COLLAPSE_WHITESPACE = re.compile(r"\s+")
_RE_COLLAPSE_EXTRA_BLANK_LINES = re.compile(r"\n{3,}")
_RE_ESCAPE_MARKDOWN_LEADING = re.compile(r"^([>#\-\+\*])")
_RE_ESCAPE_MARKDOWN_ORDERED_LIST = re.compile(r"^(\d+)\.\s")
_RE_WRAP_INLINE_MARKDOWN = re.compile(r"^(\s*)(.*?)(\s*)$", re.DOTALL)

def _configure_windows_stdio():
    """
    Windows 控制台经常使用 GBK/CP936，直接输出某些符号（如 ✓/✗）会触发 UnicodeEncodeError。

    策略：
    - 交互式控制台（isatty=True）：不强制切换编码，尽量保持用户终端显示正常；仅将 errors 调整为 replace，避免崩溃。
    - 非交互（被管道/测试框架捕获）：优先输出 UTF-8，保证机器可读；同样使用 errors=replace 兜底。
    """
    if sys.platform != "win32":
        return

    def _safe_reconfigure(stream, *, encoding=None, errors=None):
        if not hasattr(stream, "reconfigure"):
            return False
        try:
            kwargs = {}
            if encoding is not None:
                kwargs["encoding"] = encoding
            if errors is not None:
                kwargs["errors"] = errors
            stream.reconfigure(**kwargs)
            return True
        except Exception:
            return False

    def _safe_wrap(stream, *, encoding, errors):
        buffer = None
        if hasattr(stream, "detach"):
            try:
                buffer = stream.detach()
            except Exception:
                buffer = None
        if buffer is None:
            buffer = getattr(stream, "buffer", None)
        if buffer is None:
            return False
        try:
            wrapped = io.TextIOWrapper(buffer, encoding=encoding, errors=errors, line_buffering=True)
            if stream is sys.stdout:
                sys.stdout = wrapped
            elif stream is sys.stderr:
                sys.stderr = wrapped
            return True
        except Exception:
            return False

    is_tty = bool(getattr(sys.stdout, "isatty", lambda: False)())
    errors = "replace"

    if is_tty:
        if not _safe_reconfigure(sys.stdout, errors=errors):
            current_encoding = getattr(sys.stdout, "encoding", None) or "utf-8"
            _safe_wrap(sys.stdout, encoding=current_encoding, errors=errors)
        if not _safe_reconfigure(sys.stderr, errors=errors):
            current_encoding = getattr(sys.stderr, "encoding", None) or "utf-8"
            _safe_wrap(sys.stderr, encoding=current_encoding, errors=errors)
        return

    target_encoding = "utf-8"
    if not _safe_reconfigure(sys.stdout, encoding=target_encoding, errors=errors):
        _safe_wrap(sys.stdout, encoding=target_encoding, errors=errors)
    if not _safe_reconfigure(sys.stderr, encoding=target_encoding, errors=errors):
        _safe_wrap(sys.stderr, encoding=target_encoding, errors=errors)


_configure_windows_stdio()

# ==================== 依赖配置 ====================

_DEPENDENCIES_BY_EXT = {
    '.docx': [('docx', 'python-docx')],
    '.xlsx': [('openpyxl', 'openpyxl')],
    '.pptx': [('pptx', 'python-pptx')],
    '.pdf': [('pdfplumber', 'pdfplumber')],
    '.md': [],  # Markdown 转 DOCX 使用 Node.js，无 Python 依赖
}

# ==================== Node.js 共享依赖目录 ====================

def _get_node_shared_root():
    override = os.environ.get(NODE_SHARED_HOME_ENV)
    if override:
        return override

    if sys.platform == "win32":
        base = os.environ.get("LOCALAPPDATA") or os.environ.get("APPDATA")
        if not base:
            base = os.path.join(os.path.expanduser("~"), "AppData", "Local")
        return os.path.join(base, "BruceDocConverter", "node")

    return os.path.join(os.path.expanduser("~"), ".bruce-doc-converter", "node")

def _sync_shared_package_files(source_dir, target_dir):
    for filename in ("package.json", "package-lock.json"):
        src = os.path.join(source_dir, filename)
        if not os.path.exists(src):
            continue
        dst = os.path.join(target_dir, filename)
        try:
            shutil.copy2(src, dst)
        except Exception as e:
            return False, f"无法复制 {filename} 到共享目录: {str(e)}"
    return True, None

def _ensure_shared_node_modules(shared_dir, source_dir):
    npm_cmd = shutil.which('npm')
    if not npm_cmd:
        return False, "未找到 npm。请安装 Node.js（自带 npm）后重试。"

    try:
        os.makedirs(shared_dir, exist_ok=True)
    except Exception as e:
        return False, f"无法创建共享依赖目录: {str(e)}"

    ok, err = _sync_shared_package_files(source_dir, shared_dir)
    if not ok:
        return False, err

    cmd = [npm_cmd, "install", "--no-fund", "--no-audit"]
    try:
        print("[BruceDocConverter] 正在安装 Node.js 依赖到用户共享目录...", file=sys.stderr)
        result = subprocess.run(
            cmd,
            cwd=shared_dir,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
        )
        if result.returncode != 0:
            error_output = result.stderr.strip() or result.stdout.strip()
            return False, f"Node.js 依赖安装失败: {error_output}"
    except Exception as e:
        return False, f"Node.js 依赖安装失败: {str(e)}"

    return True, None

def _find_mmdc_binary(node_modules_dir):
    if not node_modules_dir:
        return None

    ext = ".cmd" if sys.platform == "win32" else ""
    candidate = os.path.join(node_modules_dir, ".bin", f"mmdc{ext}")
    if os.path.exists(candidate):
        return candidate
    return None

# ==================== 依赖安装函数 ====================

def install_dependencies(pip_packages):
    """
    自动安装缺失的依赖到用户目录（使用 --user 标志）

    使用 --user 安装的好处：
    - 不受 PEP 668 系统保护限制（适用于 macOS 从系统 Python 安装）
    - 无需虚拟环境
    - 安装到 ~/.local/ (Linux/macOS) 或 %APPDATA% (Windows)
    - 所有项目共享

    Args:
        pip_packages: 要安装的包名列表

    Returns:
        (success: bool, error_message: str or None)
    """
    if not pip_packages:
        return True, None

    # 获取当前 Python 可执行文件路径
    python_exe = sys.executable or 'python'

    # 构建安装命令
    # 首先尝试 --user 安装到用户目录，避免 PEP 668 限制
    # 如果失败，回退到 --break-system-packages（适用于使用系统 Python 的情况）
    install_methods = [
        ('--user', '用户目录'),
        ('--break-system-packages', '系统目录（绕过 PEP 668 保护）'),
    ]

    for install_flag, location_desc in install_methods:
        cmd = [python_exe, '-m', 'pip', 'install', install_flag] + pip_packages

        try:
            # 显示安装提示（仅在首次安装时）
            print(f"[BruceDocConverter] 正在安装缺失的依赖: {', '.join(pip_packages)}", file=sys.stderr)
            print(f"[BruceDocConverter] 安装位置: {location_desc}", file=sys.stderr)

            # 运行安装命令
            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                errors="replace",
            )

            if result.returncode == 0:
                print(f"[BruceDocConverter] 依赖安装成功！", file=sys.stderr)
                return True, None

            # 检查错误类型，决定是否尝试下一种方法
            error_output = result.stderr.strip() or result.stdout.strip()

            # 如果是 PEP 668 相关错误，尝试下一种方法
            if 'externally-managed-environment' in error_output or 'PEP 668' in error_output:
                continue

            # 其他错误不再重试
            # 检查是否是权限问题
            if "Permission denied" in error_output or "Access denied" in error_output:
                return False, f"权限不足。请尝试: pip install {install_flag} {' '.join(pip_packages)}"

            # 检查是否是 pip 不存在
            if "No module named pip" in error_output:
                return False, "pip 未安装。请先安装 pip: python -m ensurepip --upgrade"

            # 其他错误
            return False, f"依赖安装失败: {error_output}"

        except FileNotFoundError:
            return False, f"找不到 Python 解释器: {python_exe}"
        except Exception as e:
            return False, f"依赖安装时发生错误: {str(e)}"

    # 所有方法都失败
    return False, "依赖安装失败：所有安装方法都失败（包括 --user 和 --break-system-packages）"


# ==================== 依赖检查函数 ====================

def check_dependencies(file_ext=None, auto_install=True):
    """
    检查必需的依赖是否已安装（默认检查全部；传入 file_ext 时仅检查该格式所需）

    Args:
        file_ext: 文件扩展名（如 '.docx'），仅检查该格式所需的依赖
        auto_install: 是否自动安装缺失的依赖（默认 True）

    Returns:
        (success: bool, error_message: str or None)
    """
    # 确定需要检查的依赖
    deps = _DEPENDENCIES_BY_EXT.get(file_ext) if file_ext else [
        ('docx', 'python-docx'),
        ('openpyxl', 'openpyxl'),
        ('pptx', 'python-pptx'),
        ('pdfplumber', 'pdfplumber'),
    ]
    if deps is None:
        deps = [
            ('docx', 'python-docx'),
            ('openpyxl', 'openpyxl'),
            ('pptx', 'python-pptx'),
            ('pdfplumber', 'pdfplumber'),
        ]

    # 检查依赖是否已安装
    missing = []
    missing_pip_names = []
    for module_name, pip_name in deps:
        try:
            importlib.import_module(module_name)
        except ImportError:
            missing.append(module_name)
            missing_pip_names.append(pip_name)

    # 如果有缺失依赖且启用了自动安装
    if missing and auto_install:
        success, error = install_dependencies(missing_pip_names)
        if success:
            # 安装成功，重新检查
            still_missing = []
            for module_name, pip_name in deps:
                try:
                    importlib.import_module(module_name)
                except ImportError:
                    still_missing.append(pip_name)

            if still_missing:
                return False, f"依赖安装后仍无法加载: {', '.join(still_missing)}。请手动安装并检查 Python 环境。"

            return True, None
        else:
            # 安装失败
            return False, error

    # 有缺失但未启用自动安装
    if missing:
        return False, f"缺少依赖库: {', '.join(missing_pip_names)}。请运行: pip install --user {' '.join(missing_pip_names)}"

    return True, None

def _normalize_text(value, preserve_newlines=False):
    """规范化提取出来的文本，减少空白和换行噪声"""
    if value is None:
        return ""

    text = str(value).replace("\r\n", "\n").replace("\r", "\n")
    if preserve_newlines:
        lines = [_RE_COLLAPSE_WHITESPACE.sub(" ", line).strip() for line in text.split("\n")]
        text = "\n".join(line for line in lines if line)
        return _RE_COLLAPSE_EXTRA_BLANK_LINES.sub("\n\n", text).strip()

    text = _RE_COLLAPSE_WHITESPACE.sub(" ", text.replace("\n", " "))
    return text.strip()

def _escape_plain_markdown_text(text):
    """避免普通文本误触发 Markdown 标题、列表、引用等语法"""
    if not text:
        return ""

    escaped = text
    escaped = _RE_ESCAPE_MARKDOWN_LEADING.sub(r"\\\1", escaped)
    escaped = _RE_ESCAPE_MARKDOWN_ORDERED_LIST.sub(r"\\\1. ", escaped)
    return escaped

def _format_inline_markdown(text, *, bold=False, italic=False):
    """保留两侧空白后再包裹 Markdown 强调标记，避免单词粘连"""
    if not text:
        return ""

    if not (bold or italic):
        return text

    match = _RE_WRAP_INLINE_MARKDOWN.match(text)
    if not match:
        return text

    leading, core, trailing = match.groups()
    if not core:
        return text

    if bold and italic:
        wrapped = f"***{core}***"
    elif bold:
        wrapped = f"**{core}**"
    else:
        wrapped = f"*{core}*"
    return f"{leading}{wrapped}{trailing}"

def _compose_inline_markdown(groups):
    """将分组后的富文本片段转换为 Markdown，并避免误转义已生成的强调语法"""
    formatted_parts = []
    for index, ((bold, italic), text) in enumerate(groups):
        current_text = text
        if index == 0 and not (bold or italic):
            current_text = _escape_plain_markdown_text(current_text)
        formatted_parts.append(_format_inline_markdown(current_text, bold=bold, italic=italic))
    return "".join(formatted_parts)

def _resolve_docx_style_font_flag(style, attr_name):
    """沿样式继承链解析 tri-state 字体属性，返回 True/False/None"""
    visited = set()
    current = style
    while current is not None:
        style_id = id(current)
        if style_id in visited:
            break
        visited.add(style_id)

        font = getattr(current, "font", None)
        value = getattr(font, attr_name, None) if font is not None else None
        if value is not None:
            return bool(value)

        current = getattr(current, "base_style", None)

    return None

def _get_docx_heading_level(style):
    """根据 Word 样式解析标题层级，无法识别时返回 None"""
    if style is None:
        return None

    style_id = getattr(style, "style_id", "")
    if isinstance(style_id, str) and style_id:
        match = re.match(r"(?i)^heading(\d+)$", style_id.strip())
        if match:
            return int(match.group(1))

    style_name = getattr(style, "name", "")
    if isinstance(style_name, str) and style_name:
        match = re.match(r"^(?:Heading|标题)\s*(\d+)$", style_name.strip())
        if match:
            return int(match.group(1))

    return None

def _resolve_docx_run_font_flag(run, paragraph, attr_name, *, allow_paragraph_style=False):
    """解析 run 的实际粗体/斜体状态，支持字符样式和可选的段落样式继承"""
    direct_value = getattr(run.font, attr_name, None)
    if direct_value is not None:
        return bool(direct_value)

    char_style_value = _resolve_docx_style_font_flag(getattr(run, "style", None), attr_name)
    if char_style_value is not None:
        return char_style_value

    if allow_paragraph_style:
        paragraph_style_value = _resolve_docx_style_font_flag(getattr(paragraph, "style", None), attr_name)
        if paragraph_style_value is not None:
            return paragraph_style_value

    return False

def _validate_input_file(file_path):
    """校验并规范化输入文件路径"""
    if file_path is None:
        return None, "文件路径不能为空"

    normalized = os.path.abspath(os.path.normpath(os.path.expanduser(str(file_path))))
    if not os.path.exists(normalized):
        return None, f'文件不存在: {normalized}'
    if not os.path.isfile(normalized):
        return None, f'输入路径不是文件: {normalized}'
    return normalized, None

def _resolve_markdown_output_path(file_path, output_dir=None):
    """生成 Markdown 输出路径，并确保输出目录可用"""
    if output_dir:
        target_dir = os.path.abspath(os.path.normpath(os.path.expanduser(str(output_dir))))
    else:
        file_dir = os.path.dirname(file_path) or '.'
        target_dir = os.path.join(file_dir, 'Markdown')

    if os.path.exists(target_dir) and not os.path.isdir(target_dir):
        raise NotADirectoryError(f'输出路径不是目录: {target_dir}')

    os.makedirs(target_dir, exist_ok=True)
    output_filename = os.path.splitext(os.path.basename(file_path))[0] + '.md'
    return os.path.join(target_dir, output_filename)

def _iter_batch_input_files(directory, recursive=True, output_dir=None):
    """遍历批量转换输入文件，跳过已生成输出目录，避免重复处理"""
    normalized_directory = os.path.abspath(os.path.normpath(os.path.expanduser(str(directory))))
    custom_output_dir = None
    if output_dir:
        custom_output_dir = os.path.abspath(os.path.normpath(os.path.expanduser(str(output_dir))))

    def _should_skip_dir(dir_path):
        if custom_output_dir and os.path.normcase(dir_path) == os.path.normcase(custom_output_dir):
            return True
        return os.path.basename(dir_path) in GENERATED_OUTPUT_DIR_NAMES

    if recursive:
        for root, dirs, files in os.walk(normalized_directory):
            dirs[:] = [dir_name for dir_name in dirs if not _should_skip_dir(os.path.join(root, dir_name))]
            for file in files:
                if os.path.splitext(file)[1].lower() in SUPPORTED_EXTENSIONS:
                    yield os.path.join(root, file)
        return

    for file in os.listdir(normalized_directory):
        file_path = os.path.join(normalized_directory, file)
        if os.path.isfile(file_path) and os.path.splitext(file)[1].lower() in SUPPORTED_EXTENSIONS:
            yield file_path

def _normalize_table_cell(value):
    """将单元格内容规范化为安全的 Markdown 表格单元格文本"""
    text = _normalize_text(value)
    # Markdown 表格分隔符转义
    text = text.replace("|", "\\|")
    return text

def _table_position_has_content(value, occupied=False):
    """判断表格位置是否应保留，用于保留合并单元格占位"""
    return occupied or (value is not None and str(value).strip() != "")

def _get_docx_grid_span(tc):
    """读取 Word 表格单元格的横向跨列数"""
    tc_pr = getattr(tc, "tcPr", None)
    grid_span = getattr(tc_pr, "gridSpan", None) if tc_pr is not None else None
    raw_value = getattr(grid_span, "val", None) if grid_span is not None else None
    try:
        return max(int(raw_value), 1)
    except (TypeError, ValueError):
        return 1

def _is_docx_vertical_merge_continuation(tc):
    """判断 Word 表格单元格是否为纵向合并的续格"""
    tc_pr = getattr(tc, "tcPr", None)
    if tc_pr is None:
        return False

    v_merge = getattr(tc_pr, "vMerge", None)
    if v_merge is None:
        return False

    return getattr(v_merge, "val", None) != "restart"

def _extract_docx_table_cell_text(tc):
    """从 Word 表格 XML 单元格中提取文本，保留段落换行"""
    paragraphs = []
    try:
        root = ET.fromstring(tc.xml)
    except ET.ParseError:
        return ""

    for paragraph in root.findall('./w:p', DOCX_XML_NAMESPACES):
        text_nodes = paragraph.findall('.//w:t', DOCX_XML_NAMESPACES)
        paragraph_text = ''.join(node.text for node in text_nodes if node.text)
        if paragraph_text:
            paragraphs.append(paragraph_text)

    if paragraphs:
        return _normalize_table_cell("\n".join(paragraphs))

    text_nodes = root.findall('.//w:t', DOCX_XML_NAMESPACES)
    return _normalize_table_cell(''.join(node.text for node in text_nodes if node.text))

def _docx_attr(node, attr_name):
    """读取 WordprocessingML 命名空间属性值"""
    if node is None:
        return None
    return node.attrib.get(f'{{{DOCX_W_NS}}}{attr_name}')

def _build_docx_numbering_index(doc):
    """构建 numId -> 抽象编号定义 的索引，支持多级编号渲染"""
    try:
        root = ET.fromstring(doc.part.numbering_part.element.xml)
    except (AttributeError, ET.ParseError, TypeError):
        return {}, {}

    num_to_abstract = {}
    abstract_levels = {}

    for num_node in root.findall('.//w:num', DOCX_XML_NAMESPACES):
        num_id = _docx_attr(num_node, 'numId')
        abstract_node = num_node.find('./w:abstractNumId', DOCX_XML_NAMESPACES)
        abstract_id = _docx_attr(abstract_node, 'val')
        if num_id and abstract_id:
            num_to_abstract[str(num_id)] = str(abstract_id)

    for abstract_node in root.findall('.//w:abstractNum', DOCX_XML_NAMESPACES):
        abstract_id = _docx_attr(abstract_node, 'abstractNumId')
        if not abstract_id:
            continue

        levels = {}
        for level_node in abstract_node.findall('./w:lvl', DOCX_XML_NAMESPACES):
            ilvl_raw = _docx_attr(level_node, 'ilvl')
            try:
                ilvl = int(ilvl_raw)
            except (TypeError, ValueError):
                continue

            start_node = level_node.find('./w:start', DOCX_XML_NAMESPACES)
            num_fmt_node = level_node.find('./w:numFmt', DOCX_XML_NAMESPACES)
            lvl_text_node = level_node.find('./w:lvlText', DOCX_XML_NAMESPACES)

            try:
                start = int(_docx_attr(start_node, 'val') or 1)
            except (TypeError, ValueError):
                start = 1

            levels[ilvl] = {
                'start': start,
                'num_fmt': _docx_attr(num_fmt_node, 'val') or '',
                'lvl_text': _docx_attr(lvl_text_node, 'val') or f'%{ilvl + 1}.',
            }

        if levels:
            abstract_levels[str(abstract_id)] = levels

    return num_to_abstract, abstract_levels

def _get_docx_style_numpr(style):
    """从段落样式中读取继承的 numPr，兼容内建 List Number / List Bullet 样式"""
    if style is None or not hasattr(style, 'element'):
        return None, None

    try:
        num_id_nodes = style.element.xpath('./w:pPr/w:numPr/w:numId')
        ilvl_nodes = style.element.xpath('./w:pPr/w:numPr/w:ilvl')
    except Exception:
        return None, None

    num_id = _docx_attr(num_id_nodes[0], 'val') if num_id_nodes else None
    ilvl_raw = _docx_attr(ilvl_nodes[0], 'val') if ilvl_nodes else None

    try:
        ilvl = int(ilvl_raw) if ilvl_raw is not None else 0
    except (TypeError, ValueError):
        ilvl = 0

    return (str(num_id), ilvl) if num_id is not None else (None, None)

def _get_docx_paragraph_numpr(para):
    """获取段落实际使用的 numId / ilvl，优先段落自身，再回退样式"""
    p = getattr(para, "_p", None)
    ppr = getattr(p, "pPr", None) if p is not None else None
    num_pr = getattr(ppr, "numPr", None) if ppr is not None else None

    if num_pr is not None:
        num_id_el = getattr(num_pr, "numId", None)
        ilvl_el = getattr(num_pr, "ilvl", None)
        num_id = getattr(num_id_el, "val", None)
        ilvl_raw = getattr(ilvl_el, "val", None)
        try:
            ilvl = int(ilvl_raw) if ilvl_raw is not None else 0
        except (TypeError, ValueError):
            ilvl = 0
        if num_id is not None:
            return str(num_id), ilvl

    return _get_docx_style_numpr(getattr(para, "style", None))

def _to_roman(value):
    if value <= 0:
        return str(value)
    pairs = [
        (1000, "M"), (900, "CM"), (500, "D"), (400, "CD"),
        (100, "C"), (90, "XC"), (50, "L"), (40, "XL"),
        (10, "X"), (9, "IX"), (5, "V"), (4, "IV"), (1, "I"),
    ]
    result = []
    remaining = value
    for number, numeral in pairs:
        while remaining >= number:
            result.append(numeral)
            remaining -= number
    return ''.join(result)

def _to_alpha(value, uppercase=False):
    if value <= 0:
        return str(value)
    letters = []
    current = value
    while current > 0:
        current -= 1
        letters.append(chr((current % 26) + (65 if uppercase else 97)))
        current //= 26
    return ''.join(reversed(letters))

def _to_chinese_counting(value):
    if value <= 0:
        return str(value)

    digits = "零一二三四五六七八九"
    units = ["", "十", "百", "千"]
    if value < 10:
        return digits[value]
    if value < 10000:
        parts = []
        zero_pending = False
        chars = list(str(value))
        length = len(chars)
        for idx, char in enumerate(chars):
            digit = int(char)
            unit = units[length - idx - 1]
            if digit == 0:
                zero_pending = bool(parts)
                continue
            if zero_pending:
                parts.append("零")
                zero_pending = False
            if not (digit == 1 and unit == "十" and not parts):
                parts.append(digits[digit])
            parts.append(unit)
        return ''.join(parts) or digits[0]
    return str(value)

def _to_circled_number(value):
    if 1 <= value <= 20:
        return chr(9311 + value)
    return str(value)

def _format_docx_number_value(value, num_fmt):
    """按 Word numFmt 渲染编号文本"""
    fmt = (num_fmt or '').lower()
    if fmt in {'decimal', 'decimalfullwidth'}:
        return str(value)
    if fmt == 'decimalzero':
        return f"{value:02d}"
    if fmt == 'lowerletter':
        return _to_alpha(value, uppercase=False)
    if fmt == 'upperletter':
        return _to_alpha(value, uppercase=True)
    if fmt == 'lowerroman':
        return _to_roman(value).lower()
    if fmt == 'upperroman':
        return _to_roman(value)
    if fmt in {'chinesecounting', 'chineselegalsimplified', 'ideographtraditional', 'taiwanesecounting'}:
        return _to_chinese_counting(value)
    if fmt in {'decimalenclosedcircle', 'circleNumDbPlain'.lower(), 'decimalenclosedcirclechinese'}:
        return _to_circled_number(value)
    return str(value)

def _render_docx_list_marker(numbering_info, numbering_state):
    """渲染 Word 多级编号列表的 Markdown 前缀"""
    if not numbering_info:
        return None
    if not numbering_info["ordered"]:
        return "-"

    num_id = numbering_info["num_id"]
    level = numbering_info["level"]
    levels = numbering_info["levels"]
    level_def = levels.get(level, {})
    state = numbering_state.setdefault(num_id, {})

    for existing_level in list(state.keys()):
        if existing_level > level:
            del state[existing_level]

    if level not in state:
        state[level] = max(level_def.get("start", 1) - 1, 0)
    state[level] += 1

    for ancestor_level in range(level):
        if ancestor_level not in state:
            ancestor_def = levels.get(ancestor_level, {})
            state[ancestor_level] = max(ancestor_def.get("start", 1), 1)

    template = level_def.get("lvl_text") or f"%{level + 1}."

    def _replace(match):
        ref_level = int(match.group(1)) - 1
        ref_value = state.get(ref_level, 1)
        ref_def = levels.get(ref_level, level_def)
        return _format_docx_number_value(ref_value, ref_def.get("num_fmt"))

    rendered = re.sub(r"%(\d+)", _replace, template).strip()
    return rendered or "1."

def _is_docx_toc_paragraph(para):
    """识别 Word 自动目录段落，避免被误当正文导出"""
    style = getattr(para, "style", None)
    style_name = getattr(style, "name", "") if style is not None else ""
    style_id = getattr(style, "style_id", "") if style is not None else ""

    if re.match(r"(?i)^toc(?:\s+heading|\s+\d+)?$", style_name.strip()):
        return True
    if re.match(r"(?i)^toc(?:heading|\d+)?$", style_id.strip()):
        return True

    try:
        root = ET.fromstring(para._p.xml)
    except Exception:
        return False

    for instr in root.findall('.//w:instrText', DOCX_XML_NAMESPACES):
        if 'TOC' in (instr.text or '').upper():
            return True
    return False

# ==================== 图片提取公共基础设施 ====================

def _detect_image_format(data):
    """通过文件头魔数识别图片格式，返回扩展名（不含点）或 None"""
    if not data or len(data) < 8:
        return None
    for signature, fmt in _IMAGE_SIGNATURES.items():
        if data[:len(signature)] == signature:
            # EMF 需要额外验证：前 4 字节为 \x01\x00\x00\x00 且偏移 40 处有 ' EMF' 标记
            if fmt == 'emf' and len(data) >= 44:
                if data[40:44] != b' EMF':
                    continue
            return fmt
    return None


def _get_image_dimensions(data):
    """
    从图片二进制数据中解析宽高（像素）。
    不依赖 PIL，仅通过文件头解析。
    对于不支持解析的格式返回 (None, None)。
    """
    if not data or len(data) < 8:
        return None, None

    fmt = _detect_image_format(data)
    if fmt is None:
        return None, None

    try:
        if fmt == 'png':
            # PNG: IHDR chunk 位于文件头之后，偏移 16 处为宽高各 4 字节大端
            if len(data) >= 24:
                width = struct.unpack('>I', data[16:20])[0]
                height = struct.unpack('>I', data[20:24])[0]
                return width, height

        elif fmt == 'jpeg':
            # JPEG: 扫描 SOF marker (0xFF 0xC0-0xCF, 排除 0xC4/0xC8/0xCC)
            offset = 2
            while offset < len(data) - 9:
                if data[offset] != 0xFF:
                    break
                marker = data[offset + 1]
                if marker == 0xD9:  # EOI
                    break
                if marker == 0xDA:  # SOS - 数据流开始，停止扫描
                    break
                length = struct.unpack('>H', data[offset + 2:offset + 4])[0]
                # SOF markers: 0xC0-0xCF 但排除 DHT(0xC4)、JPG(0xC8)、DAC(0xCC)
                if 0xC0 <= marker <= 0xCF and marker not in (0xC4, 0xC8, 0xCC):
                    if offset + 9 <= len(data):
                        height = struct.unpack('>H', data[offset + 5:offset + 7])[0]
                        width = struct.unpack('>H', data[offset + 7:offset + 9])[0]
                        return width, height
                offset += 2 + length

        elif fmt == 'gif':
            # GIF: 宽高位于偏移 6 处，各 2 字节小端
            if len(data) >= 10:
                width = struct.unpack('<H', data[6:8])[0]
                height = struct.unpack('<H', data[8:10])[0]
                return width, height

        elif fmt == 'bmp':
            # BMP: 宽高位于 DIB header 中，偏移 18 处各 4 字节小端（有符号）
            if len(data) >= 26:
                width = struct.unpack('<i', data[18:22])[0]
                height = abs(struct.unpack('<i', data[22:26])[0])
                return width, height

        elif fmt == 'tiff':
            # TIFF 解析复杂，跳过尺寸检测
            pass

    except (struct.error, IndexError):
        pass

    return None, None


def _is_decorative_image(data, width=None, height=None, is_decorative_flag=False,
                         is_pptx_background=False):
    """
    综合判定图片是否为装饰性/无意义图片。

    Args:
        data: 图片二进制数据
        width: 图片宽度（像素），None 时自动检测
        height: 图片高度（像素），None 时自动检测
        is_decorative_flag: Office 文档中 adec:decorative 标记
        is_pptx_background: PowerPoint 中覆盖整个幻灯片的背景图

    Returns:
        True 表示应过滤掉此图片
    """
    # 1. Office 自身的装饰性标记（最可靠）
    if is_decorative_flag:
        return True

    # 2. PowerPoint 全屏背景图
    if is_pptx_background:
        return True

    if not data:
        return True

    # 3. 数据量极小（纯色/透明占位）
    if len(data) <= MIN_IMAGE_DATA_BYTES:
        return True

    # 自动检测尺寸
    if width is None or height is None:
        width, height = _get_image_dimensions(data)

    # 4. 尺寸过小（项目符号图标、边框像素等）
    if width is not None and height is not None:
        if width <= MIN_IMAGE_DIMENSION_PX and height <= MIN_IMAGE_DIMENSION_PX:
            return True

        # 5. 尺寸过窄/过扁（分隔线、装饰条）
        if width > 0 and height > 0:
            ratio = max(width / height, height / width)
            if ratio >= MAX_ASPECT_RATIO:
                return True

    return False


def _setup_image_output_dir(markdown_output_path):
    """
    在 Markdown 输出文件旁创建 images/ 子目录。

    Args:
        markdown_output_path: Markdown 输出文件的绝对路径

    Returns:
        (image_save_dir, image_rel_dir) 绝对路径和相对路径
    """
    md_dir = os.path.dirname(markdown_output_path)
    image_save_dir = os.path.join(md_dir, IMAGE_OUTPUT_DIR_NAME)
    os.makedirs(image_save_dir, exist_ok=True)
    return image_save_dir, IMAGE_OUTPUT_DIR_NAME


def _save_extracted_image(data, image_save_dir, image_rel_dir, base_name, image_counter):
    """
    保存图片到 images/ 目录。

    Args:
        data: 图片二进制数据
        image_save_dir: images/ 目录绝对路径
        image_rel_dir: images/ 相对路径（用于 Markdown 引用）
        base_name: 文档基础名称（不含扩展名）
        image_counter: 图片序号

    Returns:
        相对路径字符串，如 'images/report_img_001.png'，失败返回 None
    """
    fmt = _detect_image_format(data)
    if fmt is None:
        # 尝试猜测，默认保存为 png
        fmt = 'png'

    # 对 wmf/emf 保持原格式
    ext = fmt
    filename = f"{base_name}_img_{image_counter:03d}.{ext}"
    abs_path = os.path.join(image_save_dir, filename)

    try:
        with open(abs_path, 'wb') as f:
            f.write(data)
        # 使用正斜杠确保 Markdown 跨平台兼容
        return f"{image_rel_dir}/{filename}"
    except OSError:
        logger.debug("Failed to save extracted image: %s", abs_path, exc_info=True)
        return None


def _make_image_markdown(rel_path, alt_text=None):
    """生成 Markdown 图片语法"""
    alt = alt_text.strip() if alt_text else "image"
    # 转义 alt 文本中的 Markdown 特殊字符
    alt = alt.replace('[', '\\[').replace(']', '\\]')
    return f"![{alt}]({rel_path})"


def _check_ooxml_decorative_flag(element, namespaces=None):
    """
    检查 OOXML 元素是否标记为装饰性（adec:decorative val="1"）。
    同时获取元素的 alt text（descr 属性）。

    Args:
        element: lxml/ET 元素（通常是 docPr 或 cNvPr）
        namespaces: 命名空间字典

    Returns:
        (is_decorative, alt_text) 元组
    """
    if element is None:
        return False, ""

    ns = namespaces or OOXML_IMAGE_NAMESPACES
    is_decorative = False
    alt_text = ""

    # 尝试从 docPr / cNvPr 读取 descr 属性
    alt_text = element.get('descr', '') or ''

    # 检查 adec:decorative 子元素
    adec_ns = ns.get('adec', 'http://schemas.microsoft.com/office/drawing/2017/decorative')
    for child in element:
        tag = child.tag
        # 处理带命名空间和不带命名空间两种情况
        if tag == f'{{{adec_ns}}}decorative' or tag.endswith('}decorative'):
            if child.get('val', '0') == '1':
                is_decorative = True
                break

    return is_decorative, alt_text

def convert_docx(file_path, image_save_dir=None, image_rel_dir=None):
    """转换 Word 文档，支持标题、格式、列表（含编号/层级）和图片提取"""
    import docx

    doc = docx.Document(file_path)
    content = ""
    num_to_abstract, abstract_levels = _build_docx_numbering_index(doc)
    numbering_state = {}
    image_counter = 0
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    extracted_images = []

    def _extract_drawing_images(paragraph_element):
        """
        从段落 XML 中提取图片。
        支持 w:drawing（内联和浮动）以及 mc:AlternateContent 包裹的图片。

        Returns:
            图片 Markdown 字符串列表
        """
        nonlocal image_counter

        if image_save_dir is None:
            return []

        image_markdowns = []
        ns = OOXML_IMAGE_NAMESPACES

        try:
            p_xml = ET.fromstring(paragraph_element.xml)
        except (ET.ParseError, AttributeError):
            return []

        # 查找所有 drawing 元素（直接和通过 mc:AlternateContent 包裹的）
        drawings = []
        # 直接 w:drawing
        for drawing in p_xml.findall('.//w:drawing', ns):
            drawings.append(drawing)
        # mc:AlternateContent -> mc:Choice -> w:drawing
        for alt_content in p_xml.findall('.//mc:AlternateContent', ns):
            for choice in alt_content.findall('.//mc:Choice', ns):
                for drawing in choice.findall('.//w:drawing', ns):
                    if drawing not in drawings:
                        drawings.append(drawing)
            # mc:Fallback 中也可能有图片
            for fallback in alt_content.findall('.//mc:Fallback', ns):
                for drawing in fallback.findall('.//w:drawing', ns):
                    if drawing not in drawings:
                        drawings.append(drawing)

        for drawing in drawings:
            # 获取 docPr 以检查装饰性标记和 alt text
            doc_pr = drawing.find('.//wp:docPr', ns)
            if doc_pr is None:
                doc_pr = drawing.find('.//{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}docPr')

            is_decorative = False
            alt_text = ""
            if doc_pr is not None:
                is_decorative, alt_text = _check_ooxml_decorative_flag(doc_pr, ns)

            # 获取图片数据：通过 a:blip 的 r:embed 属性
            blip = drawing.find('.//a:blip', ns)
            if blip is None:
                continue

            embed_id = blip.get(f'{{{ns["r"]}}}embed') or blip.get('embed')
            if not embed_id:
                continue

            # 从 document part 的 related_parts 获取图片数据
            try:
                image_part = doc.part.related_parts.get(embed_id)
                if image_part is None:
                    continue
                image_data = image_part.blob
            except Exception:
                logger.debug("Failed to read DOCX image part: %s", embed_id, exc_info=True)
                continue

            if not image_data:
                continue

            # 装饰性过滤
            if _is_decorative_image(image_data, is_decorative_flag=is_decorative):
                continue

            # 保存图片
            image_counter += 1
            rel_path = _save_extracted_image(
                image_data, image_save_dir, image_rel_dir,
                base_name, image_counter
            )
            if rel_path:
                extracted_images.append(rel_path)
                md = _make_image_markdown(rel_path, alt_text)
                image_markdowns.append(md)

        return image_markdowns

    def get_numbering_info(para):
        """
        尝试从段落的 numPr / numbering.xml 解析列表信息

        Returns:
            None 或 {'level': int, 'ordered': bool}
        """
        try:
            num_id, level = _get_docx_paragraph_numpr(para)
            if num_id is None:
                return None

            abstract_id = num_to_abstract.get(str(num_id))
            levels = abstract_levels.get(abstract_id, {}) if abstract_id is not None else {}
            level_def = levels.get(level) or levels.get(0) or {}
            num_fmt = level_def.get('num_fmt')

            style = getattr(para, "style", None)
            style_name = getattr(style, "name", "") if style is not None else ""
            style_id = getattr(style, "style_id", "") if style is not None else ""
            style_hint = f"{style_name} {style_id}".lower()

            if (num_fmt or "").lower() == "bullet":
                ordered = False
            elif (num_fmt or ""):
                ordered = True
            elif "bullet" in style_hint or "项目符号" in style_hint or "符号" in style_hint:
                ordered = False
            elif "number" in style_hint or "编号" in style_hint:
                ordered = True
            else:
                ordered = True

            return {
                'level': max(level, 0),
                'ordered': ordered,
                'num_id': str(num_id),
                'levels': levels,
            }
        except Exception:
            return None

    def process_paragraph(para):
        """处理单个段落，识别标题、列表和格式"""
        if _is_docx_toc_paragraph(para):
            return ""
        if not para.text.strip():
            return ""

        style = para.style if hasattr(para, "style") else None
        style_name = style.name if style else ""
        style_id = getattr(style, "style_id", "") if style else ""
        heading_level = _get_docx_heading_level(style)
        allow_paragraph_style = heading_level is None

        # 先拼接富文本（列表项也需要保留粗体/斜体）
        # 将相邻同格式的 run 合并后再添加 Markdown 标记，避免 **text1****text2** 碎片
        groups = []
        for run in para.runs:
            text = run.text
            if not text:
                continue
            fmt = (
                _resolve_docx_run_font_flag(run, para, "bold", allow_paragraph_style=allow_paragraph_style),
                _resolve_docx_run_font_flag(run, para, "italic", allow_paragraph_style=allow_paragraph_style),
            )
            if groups and groups[-1][0] == fmt:
                groups[-1] = (fmt, groups[-1][1] + text)
            else:
                groups.append((fmt, text))
        formatted_text = _compose_inline_markdown(groups)
        text_value = _normalize_text(formatted_text.strip() or para.text.strip())
        if not text_value:
            return ""

        # 识别标题层级
        if heading_level is not None:
            heading_prefix = "#" * min(heading_level, 6)  # Markdown最多支持6级标题
            return f"{heading_prefix} {text_value}\n\n"

        # 检查是否是列表项
        # 优先使用 numPr + numbering.xml 解析列表编号格式与层级
        numbering_info = get_numbering_info(para)
        if numbering_info:
            indent = "    " * numbering_info["level"]
            marker = _render_docx_list_marker(numbering_info, numbering_state)
            return f"{indent}{marker} {text_value}\n"

        # 注意：python-docx对列表的支持有限，这里做基本处理（按样式兜底）
        style_name_str = style_name or ""
        style_id_str = style_id or ""
        if (isinstance(style_id_str, str) and style_id_str.startswith("List")) or (
            isinstance(style_name_str, str) and style_name_str.startswith("List")
        ):
            level = 0
            m = re.search(r"(\d+)$", style_name_str.strip())
            if m:
                level = max(int(m.group(1)) - 1, 0)
            indent = "    " * level
            if "Bullet" in style_id_str or "Bullet" in style_name_str or style_name_str == "List Bullet":
                return f"{indent}- {text_value}\n"
            if "Number" in style_id_str or "Number" in style_name_str or style_name_str == "List Number":
                return f"{indent}1. {text_value}\n"

        return text_value + "\n\n"

    # 处理文档中的所有元素（段落和表格）
    # 需要按照它们在文档中的顺序处理
    paragraphs_iter = iter(doc.paragraphs)
    tables_iter = iter(doc.tables)
    for element in doc.element.body:
        # 处理段落
        if element.tag.endswith('p'):
            para = next(paragraphs_iter, None)
            if para is not None:
                content += process_paragraph(para)
                # 提取段落中的图片
                for img_md in _extract_drawing_images(para._p):
                    content += f"\n{img_md}\n\n"

        # 处理表格
        elif element.tag.endswith('tbl'):
            table = next(tables_iter, None)
            if table is not None:
                # 使用底层 XML 读取真实网格，避免 python-docx 将合并单元格重复展开
                all_rows_data = []
                table_grid = getattr(getattr(table._tbl, "tblGrid", None), "gridCol_lst", None)
                max_cols = len(table_grid) if table_grid is not None else 0

                for tr in table._tbl.tr_lst:
                    row_data = []
                    for tc in tr.tc_lst:
                        span = _get_docx_grid_span(tc)
                        cell_text = "" if _is_docx_vertical_merge_continuation(tc) else _extract_docx_table_cell_text(tc)
                        row_data.append(cell_text)
                        if span > 1:
                            row_data.extend([""] * (span - 1))
                    all_rows_data.append(row_data)

                if not max_cols:
                    max_cols = max((len(r) for r in all_rows_data), default=0)
                if max_cols == 0:
                    continue
                for i, row_data in enumerate(all_rows_data):
                    padded = row_data + [""] * (max_cols - len(row_data))
                    content += "| " + " | ".join(padded) + " |\n"
                    if i == 0:
                        content += "| " + " | ".join(["---"] * max_cols) + " |\n"
                content += "\n"

    return content.strip(), extracted_images

def convert_xlsx(file_path, image_save_dir=None, image_rel_dir=None):
    """转换 Excel 文件，支持多表头、空白分隔区、冻结窗格、常见格式保留和图片提取"""
    import openpyxl
    from datetime import date, datetime, time

    workbook = openpyxl.load_workbook(file_path, data_only=True)
    content = ""
    image_counter = 0
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    extracted_images = []

    def _build_merge_map(worksheet):
        """构建合并单元格续格坐标集合，用于保留占位但不重复填充值"""
        merge_map = set()
        for merge_range in worksheet.merged_cells.ranges:
            for row in range(merge_range.min_row, merge_range.max_row + 1):
                for col in range(merge_range.min_col, merge_range.max_col + 1):
                    if row == merge_range.min_row and col == merge_range.min_col:
                        continue
                    merge_map.add((row, col))
        return merge_map

    def _count_number_format_decimals(number_format):
        fmt = (number_format or "").split(";")[0]
        fmt = re.sub(r'"[^"]*"', "", fmt)
        fmt = re.sub(r"\[[^\]]*\]", "", fmt)
        match = re.search(r"\.([0#]+)", fmt)
        return len(match.group(1)) if match else 0

    def _format_excel_datetime(value):
        if isinstance(value, datetime):
            if value.time() == datetime.min.time():
                return value.strftime("%Y-%m-%d")
            return value.strftime("%Y-%m-%d %H:%M:%S")
        if isinstance(value, date):
            return value.strftime("%Y-%m-%d")
        if isinstance(value, time):
            return value.strftime("%H:%M:%S")
        return _normalize_text(value)

    def _format_excel_number(value, number_format):
        decimals = _count_number_format_decimals(number_format)
        use_grouping = "," in (number_format or "")

        if "%" in (number_format or ""):
            scaled = value * 100
            formatted = f"{scaled:,.{decimals}f}" if use_grouping else f"{scaled:.{decimals}f}"
            return f"{formatted}%"

        if decimals > 0:
            return f"{value:,.{decimals}f}" if use_grouping else f"{value:.{decimals}f}"

        if isinstance(value, float) and not value.is_integer():
            text = f"{value:,}" if use_grouping else f"{value}"
            return text.rstrip("0").rstrip(".")

        integer_value = int(round(value))
        return f"{integer_value:,}" if use_grouping else str(integer_value)

    def _format_excel_cell(cell, is_merged_placeholder=False):
        value = cell.value
        if value is None:
            return "" if is_merged_placeholder else None

        if cell.is_date or isinstance(value, (datetime, date, time)):
            return _format_excel_datetime(value)

        if isinstance(value, bool):
            return "TRUE" if value else "FALSE"

        if isinstance(value, (int, float)):
            return _format_excel_number(value, cell.number_format or "")

        return _normalize_text(value)

    def _classify_excel_cell(cell, is_merged_placeholder=False):
        if is_merged_placeholder:
            return "placeholder"
        value = cell.value
        if value is None:
            return "blank"
        if cell.is_date or isinstance(value, (datetime, date, time)):
            return "date"
        if isinstance(value, bool):
            return "text"
        if isinstance(value, (int, float)):
            return "number"
        return "text"

    def _iter_table_row_groups(worksheet, merge_map):
        current_group = []
        for row in worksheet.iter_rows(values_only=False):
            cells = []
            occupied_positions = []
            display_values = []
            kinds = []

            for cell in row:
                coord = (cell.row, cell.column)
                is_merged_placeholder = coord in merge_map
                display_value = _format_excel_cell(cell, is_merged_placeholder=is_merged_placeholder)
                occupied = _table_position_has_content(display_value, occupied=is_merged_placeholder)

                cells.append(cell)
                display_values.append(display_value)
                occupied_positions.append(occupied)
                kinds.append(_classify_excel_cell(cell, is_merged_placeholder=is_merged_placeholder))

            row_record = {
                "row_index": row[0].row if row else 0,
                "cells": cells,
                "values": display_values,
                "occupied": occupied_positions,
                "kinds": kinds,
            }

            if any(occupied_positions):
                current_group.append(row_record)
            elif current_group:
                yield current_group
                current_group = []

        if current_group:
            yield current_group

    def _split_column_segments(row_group):
        width = max((len(row["occupied"]) for row in row_group), default=0)
        active_columns = [any(idx < len(row["occupied"]) and row["occupied"][idx] for row in row_group) for idx in range(width)]
        segments = []
        start = None

        for idx, is_active in enumerate(active_columns):
            if is_active and start is None:
                start = idx
            elif not is_active and start is not None:
                segments.append((start, idx))
                start = None

        if start is not None:
            segments.append((start, width))
        return segments

    def _slice_table_rows(row_group, col_start, col_end):
        sliced_rows = []
        for row in row_group:
            values = row["values"][col_start:col_end]
            occupied = row["occupied"][col_start:col_end]
            kinds = row["kinds"][col_start:col_end]
            if not any(occupied):
                continue
            sliced_rows.append({
                "row_index": row["row_index"],
                "values": values,
                "occupied": occupied,
                "kinds": kinds,
            })
        return sliced_rows

    def _profile_table_row(row_data):
        text_count = 0
        number_count = 0
        date_count = 0
        placeholder_count = 0
        non_empty_count = 0

        for value, occupied, kind in zip(row_data["values"], row_data["occupied"], row_data["kinds"]):
            if not occupied:
                continue
            if kind == "placeholder":
                placeholder_count += 1
                continue
            if value not in (None, ""):
                non_empty_count += 1
            if kind == "text":
                text_count += 1
            elif kind == "number":
                number_count += 1
            elif kind == "date":
                date_count += 1

        return {
            "text_count": text_count,
            "number_count": number_count,
            "date_count": date_count,
            "placeholder_count": placeholder_count,
            "non_empty_count": non_empty_count,
            "looks_header": non_empty_count > 0 and text_count >= (number_count + date_count),
            "looks_data": (number_count + date_count) > text_count,
        }

    def _get_freeze_header_rows(worksheet):
        freeze_panes = worksheet.freeze_panes
        if not freeze_panes:
            return 0
        if hasattr(freeze_panes, "row"):
            return max(int(freeze_panes.row) - 1, 0)
        match = re.match(r"[A-Za-z]+(\d+)", str(freeze_panes))
        if not match:
            return 0
        return max(int(match.group(1)) - 1, 0)

    def _determine_header_row_count(table_rows, freeze_header_rows):
        if not table_rows:
            return 0

        frozen_header_count = 0
        if freeze_header_rows > 0:
            for row in table_rows:
                if row["row_index"] <= freeze_header_rows:
                    frozen_header_count += 1
                else:
                    break
            if frozen_header_count > 0:
                return min(frozen_header_count, len(table_rows))

        first_profile = _profile_table_row(table_rows[0])
        if not first_profile["looks_header"]:
            return 0

        header_count = 1
        if len(table_rows) >= 3:
            second_profile = _profile_table_row(table_rows[1])
            third_profile = _profile_table_row(table_rows[2])
            if first_profile["placeholder_count"] > 0 and second_profile["looks_header"] and third_profile["looks_data"]:
                header_count = 2

        return min(header_count, len(table_rows))

    def _build_header_labels(header_rows, col_count):
        if not header_rows:
            return [f"Column {idx + 1}" for idx in range(col_count)]

        expanded_rows = []
        for row in header_rows:
            carry_text = ""
            expanded = []
            for value, occupied in zip(row["values"], row["occupied"]):
                text = _normalize_text(value)
                if text:
                    carry_text = text
                    expanded.append(text)
                elif len(header_rows) > 1 and occupied and carry_text:
                    expanded.append(carry_text)
                else:
                    expanded.append("")
            expanded_rows.append(expanded)

        headers = []
        for col_idx in range(col_count):
            parts = []
            for row in expanded_rows:
                if col_idx >= len(row):
                    continue
                part = row[col_idx].strip()
                if part and (not parts or parts[-1] != part):
                    parts.append(part)
            headers.append(" / ".join(parts) if parts else "")
        return headers

    def _render_table_block(table_rows, freeze_header_rows):
        col_count = max((len(row["values"]) for row in table_rows), default=0)
        if col_count == 0:
            return ""

        header_count = _determine_header_row_count(table_rows, freeze_header_rows)
        header_rows = table_rows[:header_count]
        headers = _build_header_labels(header_rows, col_count)
        data_rows = table_rows[header_count:] if header_count > 0 else table_rows

        lines = [
            "| " + " | ".join(_normalize_table_cell(header) for header in headers) + " |",
            "| " + " | ".join(["---"] * col_count) + " |",
        ]

        for row in data_rows:
            row_values = []
            for idx in range(col_count):
                value = row["values"][idx] if idx < len(row["values"]) else ""
                occupied = row["occupied"][idx] if idx < len(row["occupied"]) else False
                if _table_position_has_content(value, occupied):
                    row_values.append(_normalize_table_cell(value))
                else:
                    row_values.append("")
            lines.append("| " + " | ".join(row_values) + " |")

        return "\n".join(lines)

    try:
        for sheet_name in workbook.sheetnames:
            if len(workbook.sheetnames) > 1:
                content += f"## {_normalize_text(sheet_name)}\n\n"

            worksheet = workbook[sheet_name]
            merge_map = _build_merge_map(worksheet)
            freeze_header_rows = _get_freeze_header_rows(worksheet)
            table_blocks = []

            for row_group in _iter_table_row_groups(worksheet, merge_map):
                for col_start, col_end in _split_column_segments(row_group):
                    table_rows = _slice_table_rows(row_group, col_start, col_end)
                    if not table_rows:
                        continue
                    table_markdown = _render_table_block(table_rows, freeze_header_rows)
                    if table_markdown:
                        table_blocks.append(table_markdown)

            if len(table_blocks) == 1:
                content += table_blocks[0] + "\n\n"
            elif len(table_blocks) > 1:
                for idx, table_markdown in enumerate(table_blocks, 1):
                    content += f"### Table {idx}\n\n{table_markdown}\n\n"

            # 提取 worksheet 中的嵌入图片
            if image_save_dir is not None:
                try:
                    ws_images = getattr(worksheet, '_images', []) or []
                    # 按锚定行号排序
                    sorted_images = []
                    for img_obj in ws_images:
                        try:
                            anchor = getattr(img_obj, 'anchor', None)
                            anchor_from = getattr(anchor, '_from', None) if anchor else None
                            row = getattr(anchor_from, 'row', 0) if anchor_from else 0
                            col = getattr(anchor_from, 'col', 0) if anchor_from else 0
                            sorted_images.append((row, col, img_obj))
                        except Exception:
                            logger.debug("Failed to read XLSX image anchor; falling back to default order", exc_info=True)
                            sorted_images.append((0, 0, img_obj))
                    sorted_images.sort(key=lambda x: (x[0], x[1]))

                    for _row, _col, img_obj in sorted_images:
                        try:
                            # 获取图片数据
                            img_ref = getattr(img_obj, 'ref', None) or getattr(img_obj, '_data', None)
                            image_data = None
                            if img_ref is not None:
                                # openpyxl Image 对象的图片数据
                                if hasattr(img_ref, 'read'):
                                    img_ref.seek(0)
                                    image_data = img_ref.read()
                                elif isinstance(img_ref, bytes):
                                    image_data = img_ref
                            # 回退：尝试从 _data 属性读取
                            if image_data is None and hasattr(img_obj, '_data'):
                                raw = img_obj._data
                                if callable(raw):
                                    raw = raw()
                                if isinstance(raw, bytes):
                                    image_data = raw
                                elif hasattr(raw, 'read'):
                                    raw.seek(0)
                                    image_data = raw.read()

                            if not image_data:
                                continue

                            # 装饰性过滤
                            if _is_decorative_image(image_data):
                                continue

                            image_counter += 1
                            rel_path = _save_extracted_image(
                                image_data, image_save_dir, image_rel_dir,
                                base_name, image_counter
                            )
                            if rel_path:
                                extracted_images.append(rel_path)
                                content += f"{_make_image_markdown(rel_path)}\n\n"
                        except Exception:
                            logger.debug("Failed to extract an XLSX embedded image; skipping it", exc_info=True)
                            continue
                except Exception:
                    logger.debug("Failed to inspect XLSX worksheet images; continuing without them", exc_info=True)
    finally:
        workbook.close()

    return content.strip(), extracted_images

def convert_pptx(file_path, image_save_dir=None, image_rel_dir=None):
    """转换 PowerPoint 文件，提取标题、正文、表格、图表、图片和备注"""
    import pptx
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

    presentation = pptx.Presentation(file_path)
    content = ""
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    image_counter = 0
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    extracted_images = []

    def _resolve_pptx_run_font_flag(run, paragraph, attr_name, *, allow_paragraph_style=False):
        """解析 run 的实际粗体/斜体状态，支持段落默认格式回退"""
        direct_value = getattr(run.font, attr_name, None)
        if direct_value is not None:
            return bool(direct_value)

        if allow_paragraph_style:
            paragraph_value = getattr(paragraph.font, attr_name, None)
            if paragraph_value is not None:
                return bool(paragraph_value)

        return False

    def _process_text_frame(text_frame, role="body"):
        """处理文本框，保留段落层级和格式"""
        result = ""
        for para in text_frame.paragraphs:
            if not para.text.strip():
                continue

            allow_paragraph_style = role != "title"

            # 将相邻同格式的 run 合并后再添加 Markdown 标记，避免 **text1****text2** 碎片
            groups = []
            for run in para.runs:
                text = run.text
                if not text:
                    continue
                fmt = (
                    _resolve_pptx_run_font_flag(run, para, "bold", allow_paragraph_style=allow_paragraph_style),
                    _resolve_pptx_run_font_flag(run, para, "italic", allow_paragraph_style=allow_paragraph_style),
                )
                if groups and groups[-1][0] == fmt:
                    groups[-1] = (fmt, groups[-1][1] + text)
                else:
                    groups.append((fmt, text))
            formatted = _compose_inline_markdown(groups)
            text_value = _normalize_text(formatted.strip() or para.text.strip())
            if not text_value:
                continue

            if role == "title":
                result += f"### {text_value}\n\n"
                continue

            if role == "subtitle":
                result += _escape_plain_markdown_text(text_value) + "\n\n"
                continue

            # 检查列表层级
            level = para.level if para.level else 0
            if level > 0:
                indent = "  " * level
                result += f"{indent}- {text_value}\n"
            elif hasattr(para, '_pPr') and para._pPr is not None and para._pPr.find(
                './/{http://schemas.openxmlformats.org/drawingml/2006/main}buChar') is not None:
                result += f"- {text_value}\n"
            elif hasattr(para, '_pPr') and para._pPr is not None and para._pPr.find(
                './/{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum') is not None:
                result += f"1. {text_value}\n"
            else:
                result += text_value + "\n\n"

        return result

    def _iter_shapes(shapes):
        for shape in shapes:
            yield shape
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(shape.shapes)

    def _shape_is_title(shape):
        if not getattr(shape, "is_placeholder", False):
            return False

        try:
            placeholder_type = shape.placeholder_format.type
            return placeholder_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
        except Exception:
            return False

    def _get_placeholder_type(shape):
        if not getattr(shape, "is_placeholder", False):
            return None
        try:
            return shape.placeholder_format.type
        except Exception:
            return None

    def _shape_bounds(shape):
        left = getattr(shape, "left", 0)
        top = getattr(shape, "top", 0)
        width = getattr(shape, "width", 0)
        height = getattr(shape, "height", 0)
        return {
            "left": left,
            "top": top,
            "width": width,
            "height": height,
            "right": left + width,
            "bottom": top + height,
            "center_x": left + width / 2,
        }

    def _shape_sort_key(entry):
        return (entry["top"], entry["left"], entry["height"], entry["width"])

    def _render_table_markdown(table):
        all_rows_data = []
        for row in table.rows:
            row_data = []
            seen_cells = set()
            for cell in row.cells:
                cell_id = id(cell)
                if cell_id in seen_cells:
                    continue
                seen_cells.add(cell_id)
                row_data.append(_normalize_table_cell(cell.text))
            all_rows_data.append(row_data)

        max_cols = max((len(r) for r in all_rows_data), default=0)
        if max_cols == 0:
            return ""

        content_part = ""
        for idx, row_data in enumerate(all_rows_data):
            padded = row_data + [""] * (max_cols - len(row_data))
            content_part += "| " + " | ".join(padded) + " |\n"
            if idx == 0:
                content_part += "| " + " | ".join(["---"] * max_cols) + " |\n"
        return content_part.strip()

    def _render_chart_markdown(chart):
        lines = []
        chart_title = ""
        if getattr(chart, "has_title", False):
            try:
                chart_title = _normalize_text(chart.chart_title.text_frame.text)
            except Exception:
                chart_title = ""
        lines.append(f"**Chart:** {chart_title or 'Untitled chart'}")

        series_names = []
        try:
            series_names = [_normalize_text(series.name) for series in chart.series if _normalize_text(series.name)]
        except Exception:
            series_names = []
        if series_names:
            lines.append(f"Series: {', '.join(series_names)}")

        categories = []
        try:
            categories = [_normalize_text(str(category)) for category in chart.plots[0].categories if _normalize_text(str(category))]
        except Exception:
            categories = []
        if categories:
            lines.append(f"Categories: {', '.join(categories)}")

        return "\n".join(lines).strip()

    def _render_picture_markdown(caption_text=None, image_path=None, alt_text=None):
        if image_path:
            alt = alt_text or caption_text or "image"
            md = _make_image_markdown(image_path, alt)
            if caption_text:
                return f"{md}\nCaption: {caption_text}"
            return md
        if caption_text:
            return f"**Image**\nCaption: {caption_text}"
        return "**Image**"

    def _render_diagram_markdown(shape):
        text_value = ""
        if getattr(shape, "has_text_frame", False) and getattr(shape, "text", "").strip():
            text_value = _normalize_text(shape.text)
        shape_name = _normalize_text(getattr(shape, "name", "SmartArt"))
        if text_value:
            return f"**SmartArt:** {shape_name}\n{text_value}"
        return f"**SmartArt:** {shape_name}"

    def _looks_like_title_candidate(entry):
        return (
            entry["kind"] == "text"
            and entry["top"] <= slide_height * 0.22
            and entry["width"] >= slide_width * 0.35
            and len(entry.get("raw_text", "")) <= 120
        )

    def _looks_like_subtitle_candidate(entry, title_entry):
        return (
            entry["kind"] == "text"
            and entry["top"] >= title_entry["bottom"]
            and entry["top"] <= title_entry["bottom"] + slide_height * 0.18
            and entry["width"] >= slide_width * 0.25
            and entry["center_x"] >= slide_width * 0.25
            and entry["center_x"] <= slide_width * 0.75
        )

    def _is_footer_candidate(entry):
        return (
            entry["kind"] == "text"
            and entry["bottom"] >= slide_height * 0.86
            and entry["height"] <= slide_height * 0.12
        )

    def _find_picture_caption(entries, picture_entry):
        best_entry = None
        best_distance = None
        for entry in entries:
            if entry["kind"] != "text" or entry.get("role") != "body" or entry.get("consumed"):
                continue

            horizontal_overlap = min(entry["right"], picture_entry["right"]) - max(entry["left"], picture_entry["left"])
            if horizontal_overlap <= 0:
                continue

            distance = entry["top"] - picture_entry["bottom"]
            if distance < 0 or distance > slide_height * 0.08:
                continue
            if len(entry.get("raw_text", "")) > 160:
                continue

            if best_distance is None or distance < best_distance:
                best_distance = distance
                best_entry = entry
        return best_entry

    def _render_body_entries(entries):
        usable_entries = [entry for entry in entries if entry.get("markdown")]
        if not usable_entries:
            return ""

        left_entries = []
        right_entries = []
        wide_entries = []

        for entry in usable_entries:
            if entry["right"] <= slide_width * 0.48:
                left_entries.append(entry)
            elif entry["left"] >= slide_width * 0.52:
                right_entries.append(entry)
            else:
                wide_entries.append(entry)

        parts = []
        parts.extend(entry["markdown"] for entry in sorted(wide_entries, key=_shape_sort_key))

        if left_entries and right_entries:
            parts.append("#### Left Column")
            parts.extend(entry["markdown"] for entry in sorted(left_entries, key=_shape_sort_key))
            parts.append("#### Right Column")
            parts.extend(entry["markdown"] for entry in sorted(right_entries, key=_shape_sort_key))
            return "\n\n".join(part.strip() for part in parts if part and part.strip()).strip()

        ordered_entries = sorted(usable_entries, key=_shape_sort_key)
        return "\n\n".join(entry["markdown"].strip() for entry in ordered_entries if entry["markdown"].strip()).strip()

    for i, slide in enumerate(presentation.slides, 1):
        slide_parts = []
        if len(presentation.slides) > 1:
            slide_parts.append(f"## Slide {i}")

        entries = []

        for shape in _iter_shapes(slide.shapes):
            bounds = _shape_bounds(shape)
            placeholder_type = _get_placeholder_type(shape)
            entry = {
                "shape": shape,
                "kind": None,
                "role": "body",
                "markdown": "",
                "raw_text": "",
                "placeholder_type": placeholder_type,
                **bounds,
            }

            if getattr(shape, "has_table", False):
                entry["kind"] = "table"
                entry["markdown"] = _render_table_markdown(shape.table)
            elif getattr(shape, "has_chart", False):
                entry["kind"] = "chart"
                entry["markdown"] = _render_chart_markdown(shape.chart)
            elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                entry["kind"] = "picture"
                # 提取图片数据和元数据
                entry["image_path"] = None
                entry["image_alt"] = ""
                if image_save_dir is not None:
                    try:
                        image_data = shape.image.blob
                        # 检查装饰性标记：通过 shape XML 中的 cNvPr
                        is_decorative = False
                        alt_text = ""
                        try:
                            sp_xml = ET.fromstring(shape._element.xml)
                            # PPTX 中 cNvPr 可能在 p:nvPicPr/p:cNvPr 或 nvSpPr/cNvPr
                            cnv_pr = sp_xml.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
                            if cnv_pr is None:
                                cnv_pr = sp_xml.find('.//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}cNvPr')
                            if cnv_pr is None:
                                # 用更通用的查找
                                for elem in sp_xml.iter():
                                    if elem.tag.endswith('}cNvPr') or elem.tag == 'cNvPr':
                                        cnv_pr = elem
                                        break
                            if cnv_pr is not None:
                                is_decorative, alt_text = _check_ooxml_decorative_flag(cnv_pr)
                        except (AttributeError, ET.ParseError, TypeError):
                            logger.debug("Failed to parse PPTX picture metadata; continuing without decorative flag", exc_info=True)

                        # 检查是否为背景图（覆盖面积 >= 90% 幻灯片）
                        is_background = False
                        if slide_width and slide_height:
                            shape_area = entry["width"] * entry["height"]
                            slide_area = slide_width * slide_height
                            if slide_area > 0 and shape_area / slide_area >= PPTX_BACKGROUND_COVERAGE_RATIO:
                                is_background = True

                        if not _is_decorative_image(
                            image_data,
                            is_decorative_flag=is_decorative,
                            is_pptx_background=is_background
                        ):
                            image_counter += 1
                            rel_path = _save_extracted_image(
                                image_data, image_save_dir, image_rel_dir,
                                base_name, image_counter
                            )
                            if rel_path:
                                extracted_images.append(rel_path)
                                entry["image_path"] = rel_path
                                entry["image_alt"] = alt_text
                    except Exception:
                        logger.debug("Failed to extract a PPTX picture; skipping it", exc_info=True)
            elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.DIAGRAM:
                entry["kind"] = "diagram"
                entry["markdown"] = _render_diagram_markdown(shape)
            elif getattr(shape, "has_text_frame", False) and getattr(shape, "text", "").strip():
                entry["kind"] = "text"
                entry["raw_text"] = _normalize_text(shape.text, preserve_newlines=True)
                if _shape_is_title(shape):
                    entry["role"] = "title"
                elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                    entry["role"] = "subtitle"
                elif placeholder_type in (PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.SLIDE_NUMBER):
                    entry["role"] = "footer"
                entry["markdown"] = _process_text_frame(shape.text_frame, role="title" if entry["role"] == "title" else "body")
            else:
                continue

            if entry["markdown"] or entry["kind"] in {"picture"}:
                entries.append(entry)

        entries.sort(key=_shape_sort_key)

        title_entries = [entry for entry in entries if entry["role"] == "title"]
        subtitle_entries = [entry for entry in entries if entry["role"] == "subtitle"]

        if not title_entries:
            for entry in entries:
                if _looks_like_title_candidate(entry):
                    entry["role"] = "title"
                    entry["markdown"] = _process_text_frame(entry["shape"].text_frame, role="title")
                    title_entries.append(entry)
                    break

        if title_entries and not subtitle_entries:
            title_anchor = sorted(title_entries, key=_shape_sort_key)[0]
            for entry in entries:
                if entry["role"] == "body" and _looks_like_subtitle_candidate(entry, title_anchor):
                    entry["role"] = "subtitle"
                    subtitle_entries.append(entry)
                    break

        for entry in entries:
            if entry["role"] == "body" and _is_footer_candidate(entry):
                entry["role"] = "footer"

        for entry in [entry for entry in entries if entry["kind"] == "picture"]:
            caption_entry = _find_picture_caption(entries, entry)
            caption_text = caption_entry["raw_text"] if caption_entry else None
            if caption_entry is not None:
                caption_entry["consumed"] = True
            entry["markdown"] = _render_picture_markdown(
                caption_text=caption_text,
                image_path=entry.get("image_path"),
                alt_text=entry.get("image_alt")
            )

        title_entries = sorted([entry for entry in entries if entry["role"] == "title"], key=_shape_sort_key)
        subtitle_entries = sorted([entry for entry in entries if entry["role"] == "subtitle"], key=_shape_sort_key)
        footer_entries = sorted([entry for entry in entries if entry["role"] == "footer"], key=_shape_sort_key)
        body_entries = [
            entry for entry in entries
            if entry["role"] == "body" and not entry.get("consumed") and entry["kind"] in {"text", "table"}
        ]
        visual_entries = [
            entry for entry in entries
            if entry["kind"] in {"chart", "picture", "diagram"} and entry.get("markdown")
        ]

        for entry in title_entries:
            if entry["markdown"].strip():
                slide_parts.append(entry["markdown"].strip())

        if subtitle_entries:
            subtitle_body = "\n\n".join(
                _process_text_frame(entry["shape"].text_frame, role="subtitle").strip()
                for entry in subtitle_entries
                if _process_text_frame(entry["shape"].text_frame, role="subtitle").strip()
            ).strip()
            if subtitle_body:
                slide_parts.append("#### Subtitle\n\n" + subtitle_body)

        body_markdown = _render_body_entries(body_entries)
        if body_markdown:
            slide_parts.append(body_markdown)

        if visual_entries:
            visuals_body = "\n\n".join(entry["markdown"].strip() for entry in sorted(visual_entries, key=_shape_sort_key) if entry["markdown"].strip())
            if visuals_body:
                slide_parts.append("#### Visuals\n\n" + visuals_body)

        if footer_entries:
            footer_body = "\n\n".join(
                _process_text_frame(entry["shape"].text_frame, role="subtitle").strip()
                for entry in footer_entries
                if _process_text_frame(entry["shape"].text_frame, role="subtitle").strip()
            ).strip()
            if footer_body:
                slide_parts.append("#### Footer\n\n" + footer_body)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = _normalize_text(slide.notes_slide.notes_text_frame.text, preserve_newlines=True)
            if notes_text:
                slide_parts.append(f"### Notes\n\n{notes_text}")

        slide_content = "\n\n".join(part.strip() for part in slide_parts if part and part.strip()).strip()
        if slide_content:
            content += slide_content

        if i < len(presentation.slides):
            content += "---\n\n"

    return content.strip(), extracted_images

def _render_pdf_table(table_obj):
    """将 pdfplumber 表格对象渲染为 Markdown 表格字符串"""
    table = table_obj.extract()
    if not table or len(table) == 0:
        return ""
    filtered_table = []
    for row in table:
        if row and any(cell is not None and str(cell).strip() for cell in row):
            cleaned_row = [_normalize_table_cell(cell) for cell in row]
            filtered_table.append(cleaned_row)
    if not filtered_table:
        return ""
    max_cols = max(len(r) for r in filtered_table)
    normalized = [r + [""] * (max_cols - len(r)) for r in filtered_table]
    result = ""
    for idx, row in enumerate(normalized):
        result += "| " + " | ".join(row) + " |\n"
        if idx == 0:
            result += "| " + " | ".join(["---"] * len(row)) + " |\n"
    result += "\n"
    return result

# ---------- PDF 词元级文本重建辅助函数 ----------

def _group_words_into_lines(words, y_tolerance=3):
    """将词元按 top 坐标分组成行列表，每行按 x0 排序"""
    if not words:
        return []
    sorted_words = sorted(words, key=lambda w: (w['top'], w['x0']))
    lines = []
    cur_line = [sorted_words[0]]
    cur_top = sorted_words[0]['top']
    for word in sorted_words[1:]:
        if abs(word['top'] - cur_top) <= y_tolerance:
            cur_line.append(word)
        else:
            lines.append(sorted(cur_line, key=lambda w: w['x0']))
            cur_line = [word]
            cur_top = word['top']
    lines.append(sorted(cur_line, key=lambda w: w['x0']))
    return lines


def _reconstruct_line_text(line_words):
    """从同一行的词元重建文本，词间用空格分隔"""
    if not line_words:
        return ""
    return ' '.join(w['text'] for w in line_words)


def _get_body_font_size(chars):
    """计算正文字体大小（出现频率最高的字体大小，0.5pt 精度）"""
    from collections import Counter
    sizes = [round(c.get('size', 0) * 2) / 2 for c in chars if c.get('size', 0) > 0]
    if not sizes:
        return 10.0
    return float(Counter(sizes).most_common(1)[0][0])


def _get_line_avg_font_size(line_words, page_chars):
    """通过 page.chars 计算某行的平均字体大小"""
    if not line_words or not page_chars:
        return 0.0
    line_top = min(w['top'] for w in line_words)
    line_bottom = max(w['bottom'] for w in line_words)
    margin = max((line_bottom - line_top) * 0.5, 1.0)
    relevant = [
        c for c in page_chars
        if c.get('top', 0) >= line_top - margin and c.get('bottom', 0) <= line_bottom + margin
    ]
    sizes = [c.get('size', 0) for c in relevant if c.get('size', 0) > 0]
    return sum(sizes) / len(sizes) if sizes else 0.0


def _detect_column_split(page_width, words, min_side_words=15):
    """
    检测双栏布局。检查页面中央是否存在明显空白带，是则返回分割线 x 坐标，否则返回 None。
    适用于学术论文的双栏 PDF。
    """
    if not words or len(words) < min_side_words * 2:
        return None
    mid = page_width / 2
    band = page_width * 0.07  # 中央空白带宽度的一半

    center = sum(1 for w in words if w['x0'] > mid - band and w['x1'] < mid + band)
    sides = sum(1 for w in words if w['x1'] <= mid - band or w['x0'] >= mid + band)
    total = center + sides
    if total == 0:
        return None
    # 中央词占比 < 6% 且两侧词数充足，判定为双栏
    if sides >= min_side_words * 2 and center / total < 0.06:
        return mid
    return None

def _split_pdf_words_by_columns(words, page_chars, col_split, gap=5):
    """按双栏分割词元，额外保留跨栏词元，避免标题等内容丢失"""
    left_words, right_words, spanning_words = [], [], []
    left_chars, right_chars, spanning_chars = [], [], []

    left_limit = col_split + gap
    right_limit = col_split - gap

    for word in words:
        if word['x1'] <= left_limit:
            left_words.append(word)
        elif word['x0'] >= right_limit:
            right_words.append(word)
        else:
            spanning_words.append(word)

    for char in page_chars:
        char_x0 = char.get('x0', 0)
        char_x1 = char.get('x1', 0)
        if char_x1 <= left_limit:
            left_chars.append(char)
        elif char_x0 >= right_limit:
            right_chars.append(char)
        else:
            spanning_chars.append(char)

    return left_words, right_words, spanning_words, left_chars, right_chars, spanning_chars

def _split_markdown_blocks(content):
    return [block.strip() for block in re.split(r"\n\s*\n", content or "") if block.strip()]

def _parse_pdf_academic_section_block(block):
    """识别论文常见章节标题，支持标题独占或同段内联写法"""
    raw = (block or '').strip()
    plain = re.sub(r"^#{1,6}\s*", "", raw).strip()

    appendix_match = re.match(r"^(appendix|appendices|附录)\s*([A-Za-z0-9一二三四五六七八九十]*)\s*[:：\-]?\s*(.*)$", plain, re.IGNORECASE)
    if appendix_match:
        _, suffix, inline_body = appendix_match.groups()
        heading = "Appendix"
        if suffix:
            heading = f"{heading} {suffix.strip()}"
        return {"section": "appendix", "heading": heading, "inline_body": inline_body.strip()}

    patterns = [
        ("abstract", "Abstract", [r"abstract", r"摘要"]),
        ("keywords", "Keywords", [r"keywords?", r"index terms?", r"关键词"]),
        ("references", "References", [r"references?", r"bibliography", r"参考文献"]),
    ]

    for section_key, heading, aliases in patterns:
        alias_pattern = "|".join(f"(?:{alias})" for alias in aliases)
        match = re.match(rf"^(?:{alias_pattern})\s*[:：]?\s*(.*)$", plain, re.IGNORECASE)
        if not match:
            continue
        inline_body = match.group(1).strip()
        if inline_body and re.match(r"^[A-Za-z]+$", inline_body) and inline_body.lower() == plain.lower():
            inline_body = ""
        return {"section": section_key, "heading": heading, "inline_body": inline_body}

    return None

def _is_markdown_heading_block(block):
    first_line = (block or '').strip().splitlines()[0] if (block or '').strip() else ''
    return bool(re.match(r"^#{1,6}\s+\S", first_line))

def _format_pdf_keywords_block(body_blocks):
    text = " ".join(_normalize_text(block) for block in body_blocks if block.strip())
    text = re.sub(r"^(?:keywords?|index terms?|关键词)\s*[:：]\s*", "", text, flags=re.IGNORECASE)
    if not text:
        return "## Keywords"

    parts = [part.strip() for part in re.split(r"[;,；，、]\s*", text) if part.strip()]
    if not parts:
        return "## Keywords"
    return "## Keywords\n\n" + "\n".join(f"- {part}" for part in parts)

def _format_pdf_references_block(body_blocks):
    items = []
    for block in body_blocks:
        cleaned = _normalize_text(block, preserve_newlines=True)
        for line in cleaned.splitlines():
            text = re.sub(r"^(?:\[\d+\]|\d+[.)])\s*", "", line.strip())
            if text:
                items.append(text)

    if not items:
        return "## References"
    return "## References\n\n" + "\n".join(f"1. {item}" for item in items)

def _format_pdf_academic_section(section_key, heading, body_blocks):
    if section_key == "keywords":
        return _format_pdf_keywords_block(body_blocks)
    if section_key == "references":
        return _format_pdf_references_block(body_blocks)

    body = "\n\n".join(block.strip() for block in body_blocks if block.strip()).strip()
    rendered_heading = f"## {heading}"
    return f"{rendered_heading}\n\n{body}".strip()

def _postprocess_pdf_academic_sections(content):
    """将论文常见章节归一成固定 Markdown 结构"""
    blocks = _split_markdown_blocks(content)
    if not blocks:
        return content

    result_blocks = []
    current_section = None
    current_heading = None
    current_body = []

    def _flush_current():
        nonlocal current_section, current_heading, current_body
        if current_section is None:
            return
        result_blocks.append(_format_pdf_academic_section(current_section, current_heading, current_body))
        current_section = None
        current_heading = None
        current_body = []

    for block in blocks:
        parsed = _parse_pdf_academic_section_block(block)
        if parsed:
            _flush_current()
            current_section = parsed["section"]
            current_heading = parsed["heading"]
            current_body = [parsed["inline_body"]] if parsed["inline_body"] else []
            continue

        if current_section is not None and (_is_markdown_heading_block(block) or re.match(r"^##\s+Page\s+\d+\s*$", block)):
            _flush_current()
            result_blocks.append(block)
            continue

        if current_section is not None:
            current_body.append(block)
            continue

        result_blocks.append(block)

    _flush_current()
    return "\n\n".join(block for block in result_blocks if block.strip()).strip()


def _lines_to_markdown_blocks(lines, page_chars, body_size):
    """
    将行列表（每行为词元列表）转换为 (top, markdown_text) 块列表。
    - 字体明显大于正文的行识别为标题，连续标题行合并为一个标题
    - 连续普通文本行合并成段落，行间距过大时另起段落
    """
    blocks = []
    para_lines = []
    para_top = None
    prev_bottom = None
    prev_line_height = None
    heading_lines = []   # 连续标题行暂存
    heading_top = None

    def _flush_para():
        if not para_lines:
            return
        text = _normalize_text(' '.join(para_lines))
        if text:
            blocks.append((para_top, _escape_plain_markdown_text(text) + "\n\n"))

    def _flush_heading():
        if not heading_lines:
            return
        text = ' '.join(heading_lines)
        blocks.append((heading_top, f"### {text}\n\n"))

    for line_words in lines:
        if not line_words:
            continue
        line_top = min(w['top'] for w in line_words)
        line_bottom = max(w['bottom'] for w in line_words)
        line_height = max(line_bottom - line_top, 1.0)
        line_text = _reconstruct_line_text(line_words).strip()
        if not line_text:
            continue

        # 字体大小检测：比正文大 15% 以上视为标题
        line_size = _get_line_avg_font_size(line_words, page_chars)
        is_heading = line_size > 0 and body_size > 0 and line_size >= body_size * 1.15

        # 行间距检测：使用较小行高作为参考，避免大字号标题的阈值过大
        if prev_bottom is not None:
            gap = line_top - prev_bottom
            ref_height = min(line_height, prev_line_height) if prev_line_height else line_height
            large_gap = gap > ref_height * 0.8
        else:
            large_gap = False

        if is_heading:
            _flush_para()
            para_lines = []
            para_top = None
            # 连续标题行合并：若与上一标题行无大间距则追加
            if heading_lines and not large_gap:
                heading_lines.append(line_text)
            else:
                _flush_heading()
                heading_lines = [line_text]
                heading_top = line_top
        else:
            _flush_heading()
            heading_lines = []
            heading_top = None
            if large_gap:
                _flush_para()
                para_lines = []
                para_top = None
            if para_top is None:
                para_top = line_top
            para_lines.append(line_text)

        prev_bottom = line_bottom
        prev_line_height = line_height

    _flush_heading()
    _flush_para()
    return blocks


def _extract_pdf_page_blocks(page, tables):
    """
    提取单页 PDF 的文本和表格块。
    改进：
    - 使用 extract_words() 重建文本，修复 LaTeX PDF 词间空格丢失问题
    - 基于字体大小识别标题行
    - 检测双栏布局（学术论文常见），分栏提取后顺序拼接
    """
    table_bboxes = [t.bbox for t in tables] if tables else []
    blocks = []

    # 收集表格块
    for table_obj in tables:
        md = _render_pdf_table(table_obj)
        if md:
            blocks.append((table_obj.bbox[1], md))

    # 过滤掉表格区域内的对象
    filtered_page = page
    for bbox in table_bboxes:
        filtered_page = filtered_page.filter(
            lambda obj, b=bbox: not (
                obj.get("top", 0) >= b[1] and
                obj.get("bottom", 0) <= b[3] and
                obj.get("x0", 0) >= b[0] and
                obj.get("x1", 0) <= b[2]
            )
        )

    # 使用词元提取（修复空格丢失）
    # x_tolerance=2: 学术 PDF（LaTeX）词间距约 2.7pt，默认值 3 会把相邻词粘连
    try:
        words = filtered_page.extract_words(x_tolerance=2, y_tolerance=3, keep_blank_chars=False)
    except TypeError:
        # 旧版 pdfplumber 不支持 keep_blank_chars 参数
        words = filtered_page.extract_words(x_tolerance=2, y_tolerance=3)

    # 过滤旋转文字（如 arXiv 水印），保留 upright（正向）词元
    words = [w for w in words if w.get('upright', 1)]

    page_chars = filtered_page.chars

    if not words:
        # 回退到 extract_text()
        text = _normalize_text(filtered_page.extract_text(), preserve_newlines=True)
        if text:
            blocks.append((0.0, "\n".join(_escape_plain_markdown_text(line) for line in text.splitlines()) + "\n\n"))
        blocks.sort(key=lambda b: b[0])
        return blocks

    body_size = _get_body_font_size(page_chars) if page_chars else 10.0

    # 检测双栏布局
    col_split = _detect_column_split(page.width, words)

    if col_split is not None:
        # 双栏：分别处理左右两栏，并单独保留跨栏标题/摘要等元素
        left_words, right_words, spanning_words, left_chars, right_chars, spanning_chars = _split_pdf_words_by_columns(
            words,
            page_chars,
            col_split,
        )

        left_lines = _group_words_into_lines(left_words)
        right_lines = _group_words_into_lines(right_words)
        spanning_lines = _group_words_into_lines(spanning_words)

        left_blocks = _lines_to_markdown_blocks(left_lines, left_chars, body_size)
        right_blocks = _lines_to_markdown_blocks(right_lines, right_chars, body_size)
        spanning_blocks = _lines_to_markdown_blocks(spanning_lines, spanning_chars, body_size)

        left_max = max((top for top, _ in left_blocks), default=0.0)
        offset = left_max + page.height
        right_shifted = [(top + offset, content) for top, content in right_blocks]

        blocks.extend(spanning_blocks)
        blocks.extend(left_blocks)
        blocks.extend(right_shifted)
    else:
        lines = _group_words_into_lines(words)
        text_blocks = _lines_to_markdown_blocks(lines, page_chars, body_size)
        blocks.extend(text_blocks)

    blocks.sort(key=lambda b: b[0])
    return blocks

def convert_pdf(file_path):
    """转换 PDF 文件，支持文本和表格提取，按页面位置交错输出"""
    import pdfplumber

    content_parts = []
    page_errors = []
    with pdfplumber.open(file_path) as pdf:
        for page_number, page in enumerate(pdf.pages, 1):
            try:
                tables = page.find_tables()
                blocks = _extract_pdf_page_blocks(page, tables)
            except Exception as exc:
                page_errors.append((page_number, str(exc)))
                try:
                    fallback_text = _normalize_text(page.extract_text(), preserve_newlines=True)
                except Exception:
                    fallback_text = ""
                if not fallback_text:
                    continue
                blocks = [(0.0, "\n".join(_escape_plain_markdown_text(line) for line in fallback_text.splitlines()) + "\n\n")]

            if not blocks:
                continue

            page_content = "".join(block_content for _, block_content in blocks).strip()
            if not page_content:
                continue

            if len(pdf.pages) > 1:
                content_parts.append(f"## Page {page_number}\n\n{page_content}")
            else:
                content_parts.append(page_content)

    content = "\n\n".join(content_parts).strip()
    if not content and page_errors:
        page_numbers = ", ".join(str(page_number) for page_number, _ in page_errors)
        raise ValueError(f"PDF 解析失败，无法提取任何内容。异常页码: {page_numbers}")
    return _postprocess_pdf_academic_sections(content)

def convert_md(file_path, output_dir=None):
    """
    将 Markdown 文件转换为 DOCX 格式（通过 Node.js 脚本）

    Args:
        file_path: Markdown 文件路径
        output_dir: 可选的输出目录

    Returns:
        包含 'success'、'output_path' 和可选 'error' 的字典
    """
    # 检查 Node.js 是否可用
    node_cmd = shutil.which('node')
    if not node_cmd:
        return {
            'success': False,
            'error': '未找到 Node.js。Markdown 转 DOCX 需要 Node.js 环境。请安装 Node.js: https://nodejs.org/'
        }

    # 获取 Node.js 脚本路径
    script_dir = os.path.dirname(os.path.abspath(__file__))
    node_script = os.path.join(script_dir, 'md_to_docx', 'index.js')

    if not os.path.exists(node_script):
        return {
            'success': False,
            'error': f'Node.js 转换脚本不存在: {node_script}。请运行 npm install 安装依赖。'
        }

    source_dir = os.path.join(script_dir, 'md_to_docx')
    local_node_modules = os.path.join(source_dir, 'node_modules')
    shared_root = _get_node_shared_root()
    shared_dir = os.path.join(shared_root, 'md_to_docx')
    shared_node_modules = os.path.join(shared_dir, 'node_modules')

    local_mmdc = _find_mmdc_binary(local_node_modules)
    shared_mmdc = _find_mmdc_binary(shared_node_modules)

    use_shared = False
    need_shared = (not os.path.exists(local_node_modules)) or (local_mmdc is None)
    if need_shared and shared_mmdc is None:
        ok, err = _ensure_shared_node_modules(shared_dir, source_dir)
        if not ok:
            return {
                'success': False,
                'error': (
                    f"{err}\n"
                    f"可手动安装：\n"
                    f"  1) 本地安装：cd {source_dir} && npm install\n"
                    f"  2) 共享安装：cd {shared_dir} && npm install\n"
                    f"可通过环境变量 {NODE_SHARED_HOME_ENV} 指定共享目录。"
                )
            }
        shared_mmdc = _find_mmdc_binary(shared_node_modules)

    if need_shared and shared_mmdc:
        use_shared = True

    try:
        # 调用 Node.js 脚本
        cmd = [node_cmd, node_script, file_path]
        if output_dir:
            cmd.append(output_dir)

        env = os.environ.copy()
        if use_shared:
            existing = env.get("NODE_PATH")
            if existing:
                env["NODE_PATH"] = os.pathsep.join([shared_node_modules, existing])
            else:
                env["NODE_PATH"] = shared_node_modules

        mmdc_binary = local_mmdc or shared_mmdc
        if mmdc_binary:
            env["BRUCE_DOC_CONVERTER_MMDC_PATH"] = mmdc_binary

        result = subprocess.run(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=NODE_CONVERT_TIMEOUT_SECONDS,
            env=env,
        )

        # 解析输出
        try:
            stdout_text = result.stdout or ""
            output = json.loads(stdout_text)
            return output
        except json.JSONDecodeError:
            if result.returncode == 0:
                return {
                    'success': True,
                    'output_path': (result.stdout or "").strip(),
                    'message': '转换成功'
                }
            else:
                stdout_text = (result.stdout or "").strip()
                stderr_text = (result.stderr or "").strip()
                return {
                    'success': False,
                    'error': f'Node.js 脚本输出解析失败: {stdout_text}\n{stderr_text}'
                }

    except subprocess.TimeoutExpired:
        return {
            'success': False,
            'error': '转换超时（超过2分钟）'
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'调用 Node.js 脚本失败: {str(e)}'
        }

def convert_document(file_path, extract_images=True, output_dir=None):
    """
    将文档转换为 Markdown 格式

    Args:
        file_path: 文档文件路径
        extract_images: 是否提取图片（默认 True，支持 Word/Excel/PowerPoint）
        output_dir: 可选的输出目录（默认为同目录下的 Markdown/ 子目录）

    Returns:
        包含 'success'、'markdown_content'、'output_path'、可选 'extracted_images' 和 'error' 的字典
    """
    # 验证输入文件
    file_path, input_error = _validate_input_file(file_path)
    if input_error:
        return {
            'success': False,
            'error': input_error
        }

    # 检查文件大小（限制为100MB）
    try:
        file_size = os.path.getsize(file_path)
        if file_size > MAX_FILE_SIZE_BYTES:
            return {
                'success': False,
                'error': f'文件过大: {file_size / (1024*1024):.2f}MB，超过限制 {MAX_FILE_SIZE_BYTES / (1024*1024):.0f}MB'
            }
    except OSError as e:
        return {
            'success': False,
            'error': f'无法读取文件大小: {str(e)}'
        }

    # 检查文件扩展名
    file_ext = os.path.splitext(file_path)[1].lower()
    if file_ext not in SUPPORTED_EXTENSIONS:
        return {
            'success': False,
            'error': f'不支持的文件格式: {file_ext}。支持的格式: {", ".join(SUPPORTED_EXTENSIONS)}'
        }

    # Markdown 转 DOCX 使用单独的处理流程
    if file_ext == '.md':
        return convert_md(file_path, output_dir)

    # 检查依赖（按格式按需检查，避免无关依赖阻塞）
    deps_ok, error_msg = check_dependencies(file_ext)
    if not deps_ok:
        return {
            'success': False,
            'error': error_msg
        }

    try:
        # 预先确定输出路径，以便设置图片目录
        output_path = _resolve_markdown_output_path(file_path, output_dir)

        # 设置图片提取目录
        image_save_dir = None
        image_rel_dir = None
        if extract_images and file_ext in ('.docx', '.xlsx', '.pptx'):
            image_save_dir, image_rel_dir = _setup_image_output_dir(output_path)

        # 根据文件类型转换
        extracted_images = []
        if file_ext == '.docx':
            markdown_content, extracted_images = convert_docx(
                file_path, image_save_dir=image_save_dir, image_rel_dir=image_rel_dir
            )
        elif file_ext == '.xlsx':
            markdown_content, extracted_images = convert_xlsx(
                file_path, image_save_dir=image_save_dir, image_rel_dir=image_rel_dir
            )
        elif file_ext == '.pptx':
            markdown_content, extracted_images = convert_pptx(
                file_path, image_save_dir=image_save_dir, image_rel_dir=image_rel_dir
            )
        elif file_ext == '.pdf':
            markdown_content = convert_pdf(file_path)
        else:
            return {
                'success': False,
                'error': f'不支持的文件类型: {file_ext}'
            }

        warning = None
        if not markdown_content.strip():
            if file_ext == '.pdf':
                return {
                    'success': False,
                    'error': 'PDF 未提取到任何文本或表格，文件可能是扫描件、受保护文档，或仅包含图片。请先进行 OCR 或解除保护后再试。'
                }
            warning = '未提取到任何可写入的内容，原文档可能为空，或仅包含当前版本暂不支持的对象。'

        # 保存 Markdown 文件
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)

        result = {
            'success': True,
            'markdown_content': markdown_content,
            'output_path': output_path
        }
        if extracted_images:
            result['extracted_images'] = extracted_images
        if warning:
            result['warning'] = warning
        return result

    except PermissionError as e:
        return {
            'success': False,
            'error': f'权限不足: 无法读取文件或写入输出目录 - {str(e)}'
        }
    except MemoryError:
        return {
            'success': False,
            'error': '内存不足: 文件可能过大，请尝试处理较小的文件'
        }
    except FileNotFoundError as e:
        return {
            'success': False,
            'error': f'文件未找到: {str(e)}'
        }
    except OSError as e:
        return {
            'success': False,
            'error': f'系统错误: {str(e)}'
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'转换错误 ({type(e).__name__}): {str(e)}'
        }

def batch_convert(directory, recursive=True, extract_images=True, output_dir=None):
    """
    批量转换目录中的所有支持的文档

    Args:
        directory: 要扫描的目录
        recursive: 是否递归扫描子目录
        extract_images: 是否提取图片
        output_dir: 可选的输出目录

    Returns:
        转换结果列表
    """
    results = []

    normalized_directory = os.path.abspath(os.path.normpath(os.path.expanduser(str(directory))))
    if not os.path.exists(normalized_directory):
        return [{
            'file': normalized_directory,
            'result': {
                'success': False,
                'error': f'目录不存在: {normalized_directory}'
            }
        }]
    if not os.path.isdir(normalized_directory):
        return [{
            'file': normalized_directory,
            'result': {
                'success': False,
                'error': f'输入路径不是目录: {normalized_directory}'
            }
        }]

    for file_path in _iter_batch_input_files(normalized_directory, recursive=recursive, output_dir=output_dir):
        result = convert_document(file_path, extract_images, output_dir)
        results.append({
            'file': file_path,
            'result': result
        })

    return results

def main():
    if len(sys.argv) < 2:
        print('用法: python convert_document.py <file_path> [extract_images] [output_dir]')
        print('  file_path: 文档文件路径')
        print('  extract_images: true/false (默认: true，提取图片到 images/ 子目录)')
        print('  output_dir: 可选的输出目录')
        print('')
        print('支持的格式:')
        print('  - Office/PDF 转 Markdown: .docx, .xlsx, .pptx, .pdf')
        print('  - Markdown 转 Word: .md')
        print('')
        print('批量转换: python convert_document.py --batch <directory> [recursive]')
        print('  directory: 要扫描的目录')
        print('  recursive: true/false (默认: true)')
        sys.exit(1)

    # 批量转换模式
    if sys.argv[1] == '--batch':
        if len(sys.argv) < 3:
            print('错误: 批量转换需要指定目录')
            sys.exit(1)

        directory = sys.argv[2]
        recursive = sys.argv[3].lower() == 'true' if len(sys.argv) > 3 else True

        results = batch_convert(directory, recursive)

        # 输出结果统计
        success_count = sum(1 for r in results if r['result']['success'])
        total_count = len(results)

        print(json.dumps({
            'total': total_count,
            'success': success_count,
            'failed': total_count - success_count,
            'results': results
        }, ensure_ascii=False, indent=2))

        sys.exit(0 if success_count == total_count else 1)

    # 单文件转换模式
    file_path = sys.argv[1]
    extract_images = sys.argv[2].lower() == 'true' if len(sys.argv) > 2 else True
    output_dir = sys.argv[3] if len(sys.argv) > 3 else None

    result = convert_document(file_path, extract_images, output_dir)

    # 输出结果为 JSON
    print(json.dumps(result, ensure_ascii=False, indent=2))

    sys.exit(0 if result['success'] else 1)

if __name__ == '__main__':
    main()
