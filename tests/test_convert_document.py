import base64
import tempfile
import unittest
from datetime import date
from pathlib import Path
from unittest.mock import patch

import openpyxl
from docx import Document
from docx.enum.style import WD_STYLE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx import Presentation
from pptx.util import Inches

from scripts.convert_document import (
    _detect_image_format,
    _extract_pdf_page_blocks,
    _get_image_dimensions,
    _is_decorative_image,
    _postprocess_pdf_academic_sections,
    _render_docx_list_marker,
    batch_convert,
    convert_document,
)


class ConvertDocumentTests(unittest.TestCase):
    def test_convert_docx_skips_toc_paragraphs(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "toc.docx"
            output_dir = tmp_path / "out"

            document = Document()
            if "TOC 1" not in [style.name for style in document.styles]:
                document.styles.add_style("TOC 1", WD_STYLE_TYPE.PARAGRAPH)
            document.add_paragraph("目录", style="TOC Heading")
            document.add_paragraph("第一章\t1", style="TOC 1")
            document.add_paragraph("正文开始")
            document.save(docx_path)

            result = convert_document(str(docx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertNotIn("第一章\t1", result["markdown_content"])
            self.assertNotIn("目录", result["markdown_content"])
            self.assertIn("正文开始", result["markdown_content"])

    def test_render_docx_list_marker_preserves_multilevel_and_common_formats(self):
        numbering_state = {}
        levels = {
            0: {"start": 1, "num_fmt": "decimal", "lvl_text": "%1."},
            1: {"start": 1, "num_fmt": "decimal", "lvl_text": "%1.%2."},
            2: {"start": 1, "num_fmt": "decimal", "lvl_text": "%1.%2.%3."},
        }

        marker1 = _render_docx_list_marker({"ordered": True, "num_id": "n1", "level": 0, "levels": levels}, numbering_state)
        marker2 = _render_docx_list_marker({"ordered": True, "num_id": "n1", "level": 1, "levels": levels}, numbering_state)
        marker3 = _render_docx_list_marker({"ordered": True, "num_id": "n1", "level": 2, "levels": levels}, numbering_state)

        chinese_marker = _render_docx_list_marker(
            {"ordered": True, "num_id": "n2", "level": 0, "levels": {0: {"start": 1, "num_fmt": "chineseCounting", "lvl_text": "（%1）"}}},
            {},
        )
        circled_marker = _render_docx_list_marker(
            {"ordered": True, "num_id": "n3", "level": 0, "levels": {0: {"start": 1, "num_fmt": "decimalEnclosedCircle", "lvl_text": "%1"}}},
            {},
        )

        self.assertEqual("1.", marker1)
        self.assertEqual("1.1.", marker2)
        self.assertEqual("1.1.1.", marker3)
        self.assertEqual("（一）", chinese_marker)
        self.assertEqual("①", circled_marker)

    def test_convert_xlsx_keeps_merged_cells_as_single_value(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            xlsx_path = tmp_path / "merged.xlsx"
            output_dir = tmp_path / "out"

            workbook = openpyxl.Workbook()
            worksheet = workbook.active
            worksheet["A1"] = "Merged"
            worksheet.merge_cells("A1:C1")
            worksheet["A2"] = "v1"
            worksheet["B2"] = "v2"
            worksheet["C2"] = "v3"
            workbook.save(xlsx_path)
            workbook.close()

            result = convert_document(str(xlsx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("| Merged |  |  |", result["markdown_content"])
            self.assertNotIn("| Merged | Merged | Merged |", result["markdown_content"])

    def test_convert_xlsx_supports_multi_headers_freeze_panes_and_multiple_tables(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            xlsx_path = tmp_path / "report.xlsx"
            output_dir = tmp_path / "out"

            workbook = openpyxl.Workbook()
            worksheet = workbook.active
            worksheet.freeze_panes = "A3"

            worksheet["A1"] = "Region"
            worksheet["B1"] = "Metrics"
            worksheet.merge_cells("B1:D1")
            worksheet["A2"] = "Name"
            worksheet["B2"] = "Date"
            worksheet["C2"] = "Rate"
            worksheet["D2"] = "Amount"
            worksheet["A3"] = "East"
            worksheet["B3"] = date(2024, 1, 2)
            worksheet["B3"].number_format = "yyyy-mm-dd"
            worksheet["C3"] = 0.125
            worksheet["C3"].number_format = "0.0%"
            worksheet["D3"] = 12345.6
            worksheet["D3"].number_format = "#,##0.00"

            worksheet["A6"] = 100
            worksheet["B6"] = 200
            worksheet["A7"] = 300
            worksheet["B7"] = 400

            workbook.save(xlsx_path)
            workbook.close()

            result = convert_document(str(xlsx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            markdown = result["markdown_content"]
            self.assertIn("### Table 1", markdown)
            self.assertIn("| Region / Name | Metrics / Date | Metrics / Rate | Metrics / Amount |", markdown)
            self.assertIn("| East | 2024-01-02 | 12.5% | 12,345.60 |", markdown)
            self.assertIn("### Table 2", markdown)
            self.assertIn("| Column 1 | Column 2 |", markdown)
            self.assertIn("| 100 | 200 |", markdown)
            self.assertIn("| 300 | 400 |", markdown)

    def test_convert_docx_vertical_merge_continuation_renders_blank_cell(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "vertical-merge.docx"
            output_dir = tmp_path / "out"

            document = Document()
            table = document.add_table(rows=3, cols=2)
            table.cell(0, 0).text = "Top"
            table.cell(1, 0).text = "Below"
            table.cell(0, 0).merge(table.cell(1, 0))
            table.cell(0, 1).text = "A"
            table.cell(1, 1).text = "B"
            table.cell(2, 0).text = "C"
            table.cell(2, 1).text = "D"
            document.save(docx_path)

            result = convert_document(str(docx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("| Top Below | A |", result["markdown_content"])
            self.assertIn("|  | B |", result["markdown_content"])
            self.assertEqual(1, result["markdown_content"].count("| Top Below |"))

    def test_convert_pptx_allows_non_placeholder_textbox(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "textbox.pptx"
            output_dir = tmp_path / "out"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            textbox.text_frame.text = "普通文本框"
            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("普通文本框", result["markdown_content"])
            self.assertTrue(Path(result["output_path"]).exists())

    def test_convert_pptx_sorts_visual_order_and_splits_columns(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "columns.pptx"
            output_dir = tmp_path / "out"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])

            right = slide.shapes.add_textbox(Inches(5.6), Inches(2.0), Inches(2.0), Inches(0.8))
            right.text_frame.text = "右侧内容"

            footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(2.0), Inches(0.4))
            footer.text_frame.text = "页脚"

            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6.0), Inches(0.6))
            title.text_frame.text = "课程标题"

            left = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(2.0), Inches(0.8))
            left.text_frame.text = "左侧内容"

            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            markdown = result["markdown_content"]
            self.assertLess(markdown.index("### 课程标题"), markdown.index("#### Left Column"))
            self.assertIn("#### Left Column", markdown)
            self.assertIn("左侧内容", markdown)
            self.assertIn("#### Right Column", markdown)
            self.assertIn("右侧内容", markdown)
            self.assertIn("#### Footer", markdown)
            self.assertIn("页脚", markdown)

    def test_convert_pptx_extracts_chart_subtitle_and_notes(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "chart-notes.pptx"
            output_dir = tmp_path / "out"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "季度回顾"
            slide.placeholders[1].text = "销售趋势"

            chart_data = CategoryChartData()
            chart_data.categories = ["Q1", "Q2"]
            chart_data.add_series("Sales", (12, 18))
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1),
                Inches(2),
                Inches(5.5),
                Inches(3),
                chart_data,
            ).chart
            chart.has_title = True
            chart.chart_title.text_frame.text = "Revenue"
            slide.notes_slide.notes_text_frame.text = "讲解增长原因\n补充口径说明"
            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            markdown = result["markdown_content"]
            self.assertIn("### 季度回顾", markdown)
            self.assertIn("#### Subtitle", markdown)
            self.assertIn("销售趋势", markdown)
            self.assertIn("#### Visuals", markdown)
            self.assertIn("**Chart:** Revenue", markdown)
            self.assertIn("Series: Sales", markdown)
            self.assertIn("Categories: Q1, Q2", markdown)
            self.assertIn("### Notes", markdown)
            self.assertIn("讲解增长原因", markdown)

    def test_convert_pptx_groups_picture_caption(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "picture-caption.pptx"
            output_dir = tmp_path / "out"
            image_path = tmp_path / "pixel.png"

            image_path.write_bytes(base64.b64decode(
                "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+jxM0AAAAASUVORK5CYII="
            ))

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(str(image_path), Inches(1), Inches(1.5), Inches(2), Inches(2))
            caption = slide.shapes.add_textbox(Inches(1), Inches(3.65), Inches(2.8), Inches(0.5))
            caption.text_frame.text = "这是图片说明"
            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            markdown = result["markdown_content"]
            self.assertIn("#### Visuals", markdown)
            self.assertIn("**Image**", markdown)
            self.assertIn("Caption: 这是图片说明", markdown)

    def test_convert_pptx_preserves_direct_bold_at_paragraph_start(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "direct-bold.pptx"
            output_dir = tmp_path / "out"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(2.6), Inches(1))
            paragraph = textbox.text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = "AI组件检查"
            run.font.bold = True
            tail = paragraph.add_run()
            tail.text = "。"
            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("**AI组件检查**。", result["markdown_content"])
            self.assertNotIn("\\**AI组件检查**。", result["markdown_content"])

    def test_convert_pptx_preserves_bold_from_paragraph_default_font(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "paragraph-bold.pptx"
            output_dir = tmp_path / "out"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(2.6), Inches(1))
            paragraph = textbox.text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = "AI组件检查"
            paragraph.font.bold = True
            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("**AI组件检查**", result["markdown_content"])

    def test_convert_pptx_title_does_not_add_redundant_bold_markdown(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "title-bold.pptx"
            output_dir = tmp_path / "out"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            title = slide.shapes.title.text_frame.paragraphs[0]
            title.clear()
            run = title.add_run()
            run.text = "季度回顾"
            title.font.bold = True

            subtitle = slide.placeholders[1].text_frame.paragraphs[0]
            subtitle.clear()
            subtitle_run = subtitle.add_run()
            subtitle_run.text = "销售趋势"
            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("### 季度回顾", result["markdown_content"])
            self.assertNotIn("### **季度回顾**", result["markdown_content"])

    def test_convert_docx_escapes_plain_markdown_syntax(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "plain.docx"
            output_dir = tmp_path / "out"

            document = Document()
            document.add_paragraph("1. 这是正文，不是列表")
            document.add_paragraph("# 这是正文，不是标题")
            document.save(docx_path)

            result = convert_document(str(docx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("\\1. 这是正文，不是列表", result["markdown_content"])
            self.assertIn("\\# 这是正文，不是标题", result["markdown_content"])

    def test_convert_docx_preserves_bold_from_character_style(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "char-style-bold.docx"
            output_dir = tmp_path / "out"

            document = Document()
            strong_style = document.styles.add_style("StrongInline", WD_STYLE_TYPE.CHARACTER)
            strong_style.font.bold = True

            paragraph = document.add_paragraph()
            paragraph.add_run("前文 ")
            emphasis = paragraph.add_run("AI组件检查")
            emphasis.style = strong_style
            paragraph.add_run("。")
            document.save(docx_path)

            result = convert_document(str(docx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("前文 **AI组件检查**。", result["markdown_content"])

    def test_convert_docx_preserves_direct_bold_at_paragraph_start(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "direct-bold.docx"
            output_dir = tmp_path / "out"

            document = Document()
            paragraph = document.add_paragraph()
            emphasis = paragraph.add_run("AI组件检查")
            emphasis.bold = True
            paragraph.add_run("。")
            document.save(docx_path)

            result = convert_document(str(docx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("**AI组件检查**。", result["markdown_content"])
            self.assertNotIn("\\**AI组件检查**。", result["markdown_content"])

    def test_convert_docx_preserves_bold_from_non_heading_paragraph_style(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "para-style-bold.docx"
            output_dir = tmp_path / "out"

            document = Document()
            callout_style = document.styles.add_style("Callout", WD_STYLE_TYPE.PARAGRAPH)
            callout_style.font.bold = True
            document.add_paragraph("整段强调", style="Callout")
            document.save(docx_path)

            result = convert_document(str(docx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("**整段强调**", result["markdown_content"])

    def test_convert_docx_heading_style_does_not_add_redundant_bold_markdown(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "heading.docx"
            output_dir = tmp_path / "out"

            document = Document()
            document.add_heading("章节标题", level=1)
            document.save(docx_path)

            result = convert_document(str(docx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("# 章节标题", result["markdown_content"])
            self.assertNotIn("# **章节标题**", result["markdown_content"])

    def test_convert_pdf_returns_clear_error_when_no_content_extracted(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pdf_path = tmp_path / "empty.pdf"
            pdf_path.write_bytes(b"%PDF-1.4\n%%EOF\n")

            with patch("scripts.convert_document.check_dependencies", return_value=(True, None)):
                with patch("scripts.convert_document.convert_pdf", return_value=""):
                    result = convert_document(str(pdf_path))

            self.assertFalse(result["success"])
            self.assertIn("PDF 未提取到任何文本或表格", result["error"])

    def test_batch_convert_skips_generated_output_directories(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            root = Path(tmp_dir)
            (root / "source.docx").write_bytes(b"x")
            (root / "Markdown").mkdir()
            (root / "Markdown" / "generated.md").write_text("x", encoding="utf-8")
            (root / "Word").mkdir()
            (root / "Word" / "generated.docx").write_bytes(b"x")

            seen = []

            def _fake_convert(path, *_args, **_kwargs):
                seen.append(Path(path).relative_to(root).as_posix())
                return {"success": True}

            with patch("scripts.convert_document.convert_document", side_effect=_fake_convert):
                results = batch_convert(str(root), recursive=True)

            self.assertEqual(["source.docx"], seen)
            self.assertEqual(1, len(results))

    def test_extract_pdf_page_blocks_keeps_spanning_words_in_two_column_mode(self):
        words = [
            {"text": "FULLWIDTH", "x0": 10, "x1": 90, "top": 5, "bottom": 10, "upright": 1},
        ]
        for i in range(20):
            top = 20 + i * 5
            words.append({"text": f"L{i}", "x0": 5, "x1": 20, "top": top, "bottom": top + 4, "upright": 1})
            words.append({"text": f"R{i}", "x0": 80, "x1": 95, "top": top, "bottom": top + 4, "upright": 1})

        class FakePage:
            width = 100
            height = 200
            chars = []

            def __init__(self, page_words):
                self._words = page_words

            def filter(self, _predicate):
                return self

            def extract_words(self, **_kwargs):
                return list(self._words)

            def extract_text(self):
                return ""

        blocks = _extract_pdf_page_blocks(FakePage(words), tables=[])
        rendered = "".join(content for _, content in blocks)

        self.assertIn("FULLWIDTH", rendered)

    def test_postprocess_pdf_academic_sections_normalizes_common_blocks(self):
        content = """### Abstract

This is the abstract.

Keywords: alpha, beta；gamma

### References

[1] First ref

2. Second ref

### Appendix A

Proof details.
"""

        rendered = _postprocess_pdf_academic_sections(content)

        self.assertIn("## Abstract", rendered)
        self.assertIn("This is the abstract.", rendered)
        self.assertIn("## Keywords", rendered)
        self.assertIn("- alpha", rendered)
        self.assertIn("- beta", rendered)
        self.assertIn("- gamma", rendered)
        self.assertIn("## References", rendered)
        self.assertIn("1. First ref", rendered)
        self.assertIn("1. Second ref", rendered)
        self.assertIn("## Appendix A", rendered)

    # ==================== 图片提取测试 ====================

    # 1x1 透明 PNG（用于测试装饰性过滤 —— 数据量极小）
    _TINY_PNG = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+jxM0AAAAASUVORK5CYII="
    )

    # 100x80 红色 PNG（用于测试正常图片提取）
    @staticmethod
    def _make_test_png(width=100, height=80):
        """Generate a minimal valid PNG with specified dimensions."""
        import struct
        import zlib

        def _chunk(chunk_type, data):
            raw = chunk_type + data
            crc = struct.pack('>I', zlib.crc32(raw) & 0xFFFFFFFF)
            return struct.pack('>I', len(data)) + raw + crc

        header = b'\x89PNG\r\n\x1a\n'
        ihdr_data = struct.pack('>IIBBBBB', width, height, 8, 2, 0, 0, 0)  # 8-bit RGB
        ihdr = _chunk(b'IHDR', ihdr_data)

        # Create scanlines: each row is filter byte (0) + RGB pixels
        raw_data = b''
        for _ in range(height):
            raw_data += b'\x00'  # filter byte
            raw_data += b'\xff\x00\x00' * width  # red pixels

        compressed = zlib.compress(raw_data)
        idat = _chunk(b'IDAT', compressed)
        iend = _chunk(b'IEND', b'')

        return header + ihdr + idat + iend

    def test_detect_image_format_identifies_common_formats(self):
        self.assertEqual(_detect_image_format(b'\x89PNG\r\n\x1a\n' + b'\x00' * 20), 'png')
        self.assertEqual(_detect_image_format(b'\xff\xd8\xff' + b'\x00' * 20), 'jpeg')
        self.assertEqual(_detect_image_format(b'GIF89a' + b'\x00' * 20), 'gif')
        self.assertEqual(_detect_image_format(b'BM' + b'\x00' * 20), 'bmp')
        self.assertIsNone(_detect_image_format(b'\x00\x00\x00\x00'))

    def test_get_image_dimensions_reads_png_size(self):
        png_data = self._make_test_png(200, 150)
        width, height = _get_image_dimensions(png_data)
        self.assertEqual(width, 200)
        self.assertEqual(height, 150)

    def test_is_decorative_filters_tiny_data(self):
        self.assertTrue(_is_decorative_image(b'\x00' * 100))  # 数据量极小
        self.assertTrue(_is_decorative_image(b''))              # 空数据
        self.assertTrue(_is_decorative_image(None))             # None

    def test_is_decorative_filters_by_flag(self):
        png_data = self._make_test_png(200, 150)
        self.assertTrue(_is_decorative_image(png_data, is_decorative_flag=True))
        self.assertFalse(_is_decorative_image(png_data, is_decorative_flag=False))

    def test_is_decorative_filters_pptx_background(self):
        png_data = self._make_test_png(200, 150)
        self.assertTrue(_is_decorative_image(png_data, is_pptx_background=True))

    def test_is_decorative_filters_tiny_dimensions(self):
        tiny_png = self._make_test_png(10, 10)
        self.assertTrue(_is_decorative_image(tiny_png))

    def test_is_decorative_passes_normal_image(self):
        normal_png = self._make_test_png(200, 150)
        self.assertFalse(_is_decorative_image(normal_png))

    def test_convert_docx_extracts_inline_image(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "with_image.docx"
            output_dir = tmp_path / "out"

            document = Document()
            document.add_paragraph("全图前文")
            # Add inline image
            png_data = self._make_test_png(200, 150)
            img_path = tmp_path / "test_img.png"
            img_path.write_bytes(png_data)
            document.add_picture(str(img_path), width=Inches(2))
            document.add_paragraph("图片后文")
            document.save(docx_path)

            result = convert_document(str(docx_path), extract_images=True, output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("全图前文", result["markdown_content"])
            self.assertIn("图片后文", result["markdown_content"])
            # Should have extracted image reference
            self.assertIn("![image]", result["markdown_content"])
            self.assertIn("images/", result["markdown_content"])
            # Image file should exist
            self.assertIn("extracted_images", result)
            self.assertTrue(len(result["extracted_images"]) > 0)
            img_file = output_dir / result["extracted_images"][0]
            self.assertTrue(img_file.exists())

    def test_convert_docx_no_images_when_extract_false(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "with_image.docx"
            output_dir = tmp_path / "out"

            document = Document()
            png_data = self._make_test_png(200, 150)
            img_path = tmp_path / "test_img.png"
            img_path.write_bytes(png_data)
            document.add_picture(str(img_path), width=Inches(2))
            document.save(docx_path)

            result = convert_document(str(docx_path), extract_images=False, output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertNotIn("![image]", result["markdown_content"])
            self.assertNotIn("extracted_images", result)

    def test_convert_pptx_extracts_picture_image(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "with_pic.pptx"
            output_dir = tmp_path / "out"

            png_data = self._make_test_png(300, 200)
            img_path = tmp_path / "slide_img.png"
            img_path.write_bytes(png_data)

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(str(img_path), Inches(1), Inches(1.5), Inches(3), Inches(2))
            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), extract_images=True, output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            markdown = result["markdown_content"]
            # Should have image link instead of **Image** placeholder
            self.assertIn("![", markdown)
            self.assertIn("images/", markdown)
            self.assertNotIn("**Image**", markdown)
            self.assertIn("extracted_images", result)
            self.assertTrue(len(result["extracted_images"]) > 0)

    def test_convert_pptx_filters_background_image(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "background.pptx"
            output_dir = tmp_path / "out"

            png_data = self._make_test_png(300, 200)
            img_path = tmp_path / "bg_img.png"
            img_path.write_bytes(png_data)

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            # Add full-slide-size picture (should be filtered as background)
            slide_w = presentation.slide_width
            slide_h = presentation.slide_height
            slide.shapes.add_picture(str(img_path), 0, 0, slide_w, slide_h)
            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), extract_images=True, output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            # Background image should be filtered out
            self.assertNotIn("extracted_images", result)
            self.assertNotIn("![image]", result["markdown_content"])


if __name__ == "__main__":
    unittest.main()
